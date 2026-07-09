import os
import math
import json
import pathlib
import re
import shutil
import tempfile
import time
from collections import defaultdict
from typing import Any

from PyQt6.QtWidgets import QFileDialog

from src.utils.knowledge_store import _extract_text, _sniff_csv_encoding


class AnalysisReport(dict):
    """Dict-like analysis result that also iterates over `results` for LLM robustness."""

    def _results(self):
        results = self.get("results", [])
        return results if isinstance(results, list) else []

    def __iter__(self):
        return iter(self._results())

    def __len__(self):
        return len(self._results())

    def __getitem__(self, key):
        if isinstance(key, (int, slice)):
            return self._results()[key]
        return super().__getitem__(key)


class AgentToolbox:
    MAX_TEXT_CHARS = 20000
    MAX_TEXT_FILE_SIZE = 5 * 1024 * 1024
    DEFAULT_DOMINANT_FREQUENCIES = 5

    def __init__(self, parent_widget, waapi_client):
        self.parent_widget = parent_widget
        self.waapi_client = waapi_client
        self._authorized_files = {}
        self._project_known_files = set()
        self._analysis_progress_callback = None
        self._analysis_finished_callback = None
        # Optional sink to defer destructive audio writes through the GUI's
        # confirmation pipeline. Signature: ``(path, apply_callable) -> dict``
        # (e.g. CodeExecutor.stage_audio_write). When set, normalize/batch
        # writes are staged for user confirmation instead of hitting disk
        # immediately. When None, writes apply directly (atomically).
        self.file_write_stager = None

    def set_analysis_progress_callbacks(self, progress_callback=None, finished_callback=None):
        self._analysis_progress_callback = progress_callback
        self._analysis_finished_callback = finished_callback

    def _report_analysis_progress(self, current: int, total: int, label: str = ""):
        callback = self._analysis_progress_callback
        if callable(callback):
            try:
                callback(int(current), int(total), (label or "").strip())
            except Exception:
                pass

    def _finish_analysis_progress(self):
        callback = self._analysis_finished_callback
        if callable(callback):
            try:
                callback()
            except Exception:
                pass

    @staticmethod
    def _coerce_path_value(path) -> str:
        if isinstance(path, dict):
            candidate = path.get("path") or path.get("file") or path.get("name")
            if isinstance(candidate, str) and candidate.strip():
                return candidate.strip()
            raise ValueError("Path dictionary does not contain a usable 'path' value.")
        if isinstance(path, pathlib.Path):
            return str(path)
        if isinstance(path, str):
            stripped = path.strip()
            if stripped:
                return stripped
        raise ValueError(f"Unsupported path input: {type(path).__name__}")

    @classmethod
    def _normalize_path(cls, path) -> str:
        return str(pathlib.Path(cls._coerce_path_value(path)).expanduser().resolve())

    @staticmethod
    def _is_file_ready(path: str) -> tuple[bool, str]:
        if not os.path.isfile(path):
            return False, "文件不存在。"
        try:
            size = os.path.getsize(path)
        except OSError as exc:
            return False, f"无法读取文件大小：{exc}"
        if size <= 0:
            return False, "文件大小为 0，可能仍在渲染。"
        try:
            with open(path, "rb") as handle:
                handle.read(1)
        except OSError as exc:
            return False, f"文件暂时不可读：{exc}"
        return True, ""

    @classmethod
    def _wait_for_files_ready(cls, paths: list[str], timeout: float = 30.0, interval: float = 0.25) -> tuple[list[str], list[dict[str, str]]]:
        deadline = time.monotonic() + max(float(timeout or 0), 0.0)
        pending = {path: {"size": None, "mtime": None, "stable": 0, "reason": ""} for path in paths}
        ready: list[str] = []
        failed: list[dict[str, str]] = []
        interval = max(float(interval or 0.25), 0.05)
        while pending:
            for path in list(pending):
                ok, reason = cls._is_file_ready(path)
                if not ok:
                    pending[path]["stable"] = 0
                    pending[path]["reason"] = reason
                    continue
                try:
                    stat = os.stat(path)
                except OSError as exc:
                    pending[path]["stable"] = 0
                    pending[path]["reason"] = f"无法读取文件状态：{exc}"
                    continue
                state = pending[path]
                signature = (stat.st_size, stat.st_mtime_ns)
                if signature == (state["size"], state["mtime"]):
                    state["stable"] = int(state["stable"] or 0) + 1
                else:
                    state["size"] = stat.st_size
                    state["mtime"] = stat.st_mtime_ns
                    state["stable"] = 0
                state["reason"] = "文件仍在变化，等待渲染写入完成。"
                if int(state["stable"] or 0) >= 2:
                    ready.append(path)
                    pending.pop(path, None)
            if not pending:
                break
            if time.monotonic() >= deadline:
                for path, state in pending.items():
                    failed.append({"path": path, "reason": str(state.get("reason") or "等待文件稳定超时。")})
                break
            time.sleep(interval)
        return ready, failed

    @staticmethod
    def _chunked(items: list[Any], size: int) -> list[list[Any]]:
        size = max(int(size or len(items) or 1), 1)
        return [items[index:index + size] for index in range(0, len(items), size)]

    def _remember_authorized_file(self, path: str, source: str = "user") -> dict[str, Any]:
        normalized = self._normalize_path(path)
        file_info = {
            "path": normalized,
            "name": os.path.basename(normalized),
            "size": os.path.getsize(normalized) if os.path.exists(normalized) else 0,
            "source": source,
        }
        self._authorized_files[normalized] = file_info
        return file_info

    def _remember_project_file(self, path: str) -> dict[str, Any]:
        info = self._remember_authorized_file(path, source="wwise-project")
        self._project_known_files.add(info["path"])
        return info

    def _make_project_source_entry(
        self,
        file_path: str,
        *,
        object_name: str = "",
        object_id: str = "",
        object_path: str = "",
        object_type: str = "",
        source_object_type: str = "",
    ) -> dict[str, Any]:
        remembered = self._remember_project_file(file_path)
        normalized_path = remembered["path"]
        filename = os.path.basename(normalized_path)
        return {
            **remembered,
            "name": filename,
            "file": filename,
            "path": normalized_path,
            "originalFilePath": normalized_path,
            "objectName": object_name or filename,
            "objectId": object_id,
            "objectPath": object_path,
            "objectType": object_type,
            "sourceObjectType": source_object_type,
            "type": source_object_type or object_type or "SourceFile",
            "exists": os.path.isfile(normalized_path),
        }

    def _resolve_source_control_path(self, raw_path: str) -> str:
        candidate = (raw_path or "").strip()
        if not candidate:
            return ""

        if os.path.isabs(candidate):
            return self._normalize_path(candidate)

        project_dir = ""
        if hasattr(self.waapi_client, "get_project_directory"):
            try:
                project_dir = (self.waapi_client.get_project_directory() or "").strip()
            except Exception:
                project_dir = ""

        if project_dir:
            full_path = os.path.join(project_dir, "Originals", candidate)
            return self._normalize_path(full_path)

        return self._normalize_path(candidate)

    def _resolve_project_source_path(self, raw_path: str) -> str:
        candidate = (raw_path or "").strip()
        if not candidate:
            return ""
        try:
            return self._resolve_source_control_path(candidate)
        except Exception:
            return ""

    def _list_source_control_project_entries(self) -> list[dict[str, Any]]:
        source_file_items = []
        if hasattr(self.waapi_client, "list_source_files"):
            try:
                source_file_items = self.waapi_client.list_source_files(
                    filter_mode="all",
                    recursive=True,
                    return_fields=["Path", "FileId", "Db"],
                ) or []
            except Exception:
                source_file_items = []

        entries = []
        seen_paths = set()
        for item in source_file_items:
            if not isinstance(item, dict):
                continue
            raw_path = item.get("Path") or item.get("path") or ""
            resolved_path = self._resolve_source_control_path(raw_path)
            if not resolved_path or not os.path.isfile(resolved_path) or resolved_path in seen_paths:
                continue
            seen_paths.add(resolved_path)
            entries.append(self._make_project_source_entry(resolved_path))
        return entries

    @staticmethod
    def _source_name_candidates(value: str) -> set[str]:
        text = (value or "").strip()
        if not text:
            return set()
        normalized = text.replace("\\", "/").split("/")[-1].strip()
        if not normalized:
            return set()
        stem, _ext = os.path.splitext(normalized)
        candidates = {normalized.casefold()}
        if stem:
            candidates.add(stem.casefold())
        return candidates

    @staticmethod
    def _wwise_object_name_from_file(path: str) -> str:
        stem = os.path.splitext(os.path.basename(str(path or "")))[0].strip()
        cleaned = "".join("_" if char in "\\/:*?\"<>|" else char for char in stem).strip()
        return cleaned or "ImportedAudio"

    def _match_objects_to_source_entries(self, objects: list[dict[str, Any]], source_entries: list[dict[str, Any]]):
        exact_map = defaultdict(list)
        stem_map = defaultdict(list)

        for entry in source_entries:
            if not isinstance(entry, dict):
                continue
            entry_name = (entry.get("name") or entry.get("file") or entry.get("path") or "").strip()
            for candidate in self._source_name_candidates(entry_name):
                exact_map[candidate].append(entry)
            entry_path = (entry.get("path") or "").strip()
            for candidate in self._source_name_candidates(entry_path):
                stem_map[candidate].append(entry)

        matched = []
        seen_paths = set()
        warnings = []
        for obj in objects:
            if not isinstance(obj, dict):
                continue
            tokens = self._source_name_candidates(obj.get("name") or "") | self._source_name_candidates(obj.get("path") or "")
            candidates = []
            for token in tokens:
                candidates.extend(exact_map.get(token, []))
                candidates.extend(stem_map.get(token, []))

            deduped = []
            candidate_paths = set()
            for entry in candidates:
                entry_path = entry.get("path")
                if entry_path and entry_path not in candidate_paths:
                    candidate_paths.add(entry_path)
                    deduped.append(entry)

            if len(deduped) != 1:
                object_name = (obj.get("name") or obj.get("id") or "unknown").strip()
                warnings.append(f"{object_name}: 无法从源文件列表中唯一匹配本地文件。")
                continue

            merged = dict(deduped[0])
            merged["objectName"] = (obj.get("name") or merged.get("objectName") or "").strip()
            merged["objectId"] = (obj.get("id") or merged.get("objectId") or "").strip()
            merged["objectPath"] = (obj.get("path") or merged.get("objectPath") or "").strip()
            merged_path = merged.get("path")
            if merged_path and merged_path not in seen_paths:
                seen_paths.add(merged_path)
                matched.append(merged)

        return matched, warnings

    def _get_selected_wwise_objects(self):
        selected = self.waapi_client.get_selected_objects()
        if not isinstance(selected, dict):
            return [], ["无法读取当前 Wwise 选中对象。"]
        objects = [item for item in (selected.get("objects", []) or []) if isinstance(item, dict) and item.get("id")]
        if not objects:
            return [], ["Wwise 中没有可用的选中对象。"]
        return objects, []

    def _query_descendant_objects_by_ids(self, object_ids: list[str], return_fields: list[str]):
        if not object_ids:
            return [], []

        result = self.waapi_client.call(
            "ak.wwise.core.object.get",
            args={
                "from": {"id": object_ids},
                "transform": [{"select": ["descendants"]}],
            },
            options={"return": return_fields},
        )
        error_text = self._extract_waapi_error(result)
        if error_text:
            return [], [f"查询选中对象后代失败: {error_text}"]
        items = result.get("return", []) if isinstance(result, dict) else []
        if not isinstance(items, list):
            return [], ["选中对象后代返回格式无效。"]
        return [item for item in items if isinstance(item, dict)], []

    def _get_selected_and_descendant_objects(self):
        selected_objects, warnings = self._get_selected_wwise_objects()
        if not selected_objects:
            return [], warnings

        selected_ids = [item.get("id") for item in selected_objects if item.get("id")]
        descendants, desc_warnings = self._query_descendant_objects_by_ids(
            selected_ids,
            ["id", "name", "type", "path"],
        )
        warnings.extend(desc_warnings)

        merged = []
        seen_ids = set()
        for item in [*selected_objects, *descendants]:
            if not isinstance(item, dict):
                continue
            object_id = (item.get("id") or "").strip()
            if not object_id or object_id in seen_ids:
                continue
            seen_ids.add(object_id)
            merged.append(item)
        return merged, warnings

    def _query_audio_source_parent_refs(self, audio_source_ids: list[str]):
        return self._query_objects_by_ids_resilient(audio_source_ids, ["id", "name", "type", "path", "parent"])

    def _query_audio_source_paths(self, audio_source_ids: list[str]):
        if not audio_source_ids:
            return [], []

        requested_fields = ["id", "name", "type", "path", "originalFilePath"]
        result = self.waapi_client.call(
            "ak.wwise.core.object.get",
            args={"from": {"id": audio_source_ids}},
            options={"return": requested_fields},
        )
        error_text = self._extract_waapi_error(result)
        if error_text:
            recovered_items, warnings = self._query_objects_by_ids_resilient(audio_source_ids, ["id", "name", "type", "path"])
            return recovered_items, [f"查询音频源原始路径失败: {error_text}", *warnings]

        items = result.get("return", []) if isinstance(result, dict) else []
        if not isinstance(items, list):
            return [], ["音频源原始路径返回格式无效。"]
        return [item for item in items if isinstance(item, dict)], []

    def request_user_file_access(self, title: str = "选择要授权给 Agent 读取的文件", file_filter: str = "All Files (*)", multiple: bool = True):
        if multiple:
            file_paths, _ = QFileDialog.getOpenFileNames(self.parent_widget, title, "", file_filter)
        else:
            file_path, _ = QFileDialog.getOpenFileName(self.parent_widget, title, "", file_filter)
            file_paths = [file_path] if file_path else []

        granted = []
        for path in file_paths:
            if path and os.path.isfile(path):
                granted.append(self._remember_authorized_file(path, source="user"))
        return granted

    def remember_paths(self, paths, source: str = "local"):
        remembered = []
        for path in paths or []:
            if not path:
                continue
            normalized = self._normalize_path(path)
            info = {
                "path": normalized,
                "name": os.path.basename(normalized) or normalized,
                "size": os.path.getsize(normalized) if os.path.isfile(normalized) else 0,
                "source": source,
                "is_dir": os.path.isdir(normalized),
            }
            self._authorized_files[normalized] = info
            remembered.append(info)
        return remembered

    def list_authorized_files(self):
        return list(self._authorized_files.values())

    def is_file_authorized(self, path: str) -> bool:
        try:
            normalized = self._normalize_path(path)
        except Exception:
            return False
        return normalized in self._authorized_files

    def read_user_file(self, path: str, max_chars: int = MAX_TEXT_CHARS):
        normalized = self._normalize_path(path)
        if not os.path.exists(normalized):
            raise FileNotFoundError(normalized)
        if os.path.isdir(normalized):
            raise IsADirectoryError(normalized)
        if os.path.getsize(normalized) > self.MAX_TEXT_FILE_SIZE:
            raise ValueError(f"File too large for text reading (> {self.MAX_TEXT_FILE_SIZE} bytes)")

        self.remember_paths([normalized], source="local")

        text = _extract_text(normalized)
        if not isinstance(text, str):
            text = str(text)
        if len(text) > max_chars:
            return text[:max_chars] + f"\n\n[Truncated to {max_chars} chars]"
        return text

    # ------------------------------------------------------------------
    # Structured document readers (CSV / XLSX / DOCX / PPTX)
    # ------------------------------------------------------------------

    @staticmethod
    def _truncate_rows(rows, max_rows: int):
        rows = list(rows or [])
        if len(rows) > max_rows:
            return rows[:max_rows], True
        return rows, False

    def _prepare_struct_read(self, path: str) -> str:
        """Common pre-flight: normalize, existence, size guard, authorize."""
        normalized = self._normalize_path(path)
        if not os.path.exists(normalized):
            raise FileNotFoundError(normalized)
        if os.path.isdir(normalized):
            raise IsADirectoryError(normalized)
        if os.path.getsize(normalized) > self.MAX_TEXT_FILE_SIZE:
            raise ValueError(f"File too large for structured reading (> {self.MAX_TEXT_FILE_SIZE} bytes)")
        self.remember_paths([normalized], source="local")
        return normalized

    def read_csv(self, path: str, *, max_rows: int = 10000,
                 encoding: str | None = None, delimiter: str | None = None) -> dict:
        """Read a CSV file into a structured dict.

        Returns:
            {
              "path", "encoding", "delimiter",
              "row_count", "rows": [[cell,...]], "columns": [...],
              "truncated": bool, "dtypes": {col: "int|float|str|bool|date|null"},
            }
        On error: {"error": "...", "path": ...}.
        """
        try:
            normalized = self._prepare_struct_read(path)
        except Exception as exc:
            return {"error": str(exc), "path": str(path)}

        try:
            enc = (encoding or "").strip() or _sniff_csv_encoding(normalized)

            sep = delimiter
            if sep is None:
                try:
                    import csv as _csv
                    with open(normalized, "r", encoding=enc, errors="replace") as f:
                        sample = f.read(8192)
                    if sample:
                        try:
                            dialect = _csv.Sniffer().sniff(sample, delimiters=",;\t|")
                            sep = dialect.delimiter
                        except _csv.Error:
                            sep = ","
                    else:
                        sep = ","
                except Exception:
                    sep = ","

            import pandas as pd  # type: ignore
            df = pd.read_csv(normalized, encoding=enc, sep=sep,
                             nrows=max_rows + 1, dtype=object, keep_default_na=False)
            truncated = len(df) > max_rows
            if truncated:
                df = df.iloc[:max_rows]

            # Re-infer dtypes column-by-column based on parsed strings
            inferred = {}
            for col in df.columns:
                series = df[col]
                sample_vals = [v for v in series.head(50).tolist() if v not in ("", None)]
                if not sample_vals:
                    inferred[str(col)] = "null"
                    continue
                kind = self._classify_cell_kind(sample_vals)
                inferred[str(col)] = kind

            columns = [str(c) for c in df.columns]
            rows = df.values.tolist()
            return {
                "path": normalized,
                "encoding": enc,
                "delimiter": sep,
                "row_count": int(len(rows)),
                "rows": rows,
                "columns": columns,
                "truncated": bool(truncated),
                "dtypes": inferred,
            }
        except Exception as exc:
            return {"error": f"{type(exc).__name__}: {exc}", "path": normalized}

    @staticmethod
    def _classify_cell_kind(values: list) -> str:
        all_int = True
        all_float = True
        all_bool = True
        for v in values:
            s = str(v).strip()
            if s.lower() not in ("true", "false", "0", "1", "yes", "no"):
                all_bool = False
            if not (s.lstrip("-").isdigit()):
                all_int = False
            try:
                float(s)
            except (TypeError, ValueError):
                all_float = False
        if all_int:
            return "int"
        if all_float:
            return "float"
        if all_bool:
            return "bool"
        return "str"

    def read_xlsx(self, path: str, *, sheet=None, max_rows: int = 10000) -> dict:
        """Read an XLSX workbook into a structured dict.

        sheet=None  → return all sheets (each up to max_rows)
        sheet=<str|int> → only that sheet
        """
        try:
            normalized = self._prepare_struct_read(path)
        except Exception as exc:
            return {"error": str(exc), "path": str(path)}

        try:
            import openpyxl  # type: ignore
            wb = openpyxl.load_workbook(normalized, read_only=True, data_only=True)
            try:
                wanted_titles: list[str]
                all_titles = list(wb.sheetnames)
                if sheet is None:
                    wanted_titles = all_titles
                elif isinstance(sheet, int):
                    if 0 <= sheet < len(all_titles):
                        wanted_titles = [all_titles[sheet]]
                    else:
                        return {"error": f"sheet index out of range: {sheet}", "path": normalized}
                else:
                    title = str(sheet)
                    if title not in all_titles:
                        return {"error": f"sheet not found: {title}", "path": normalized}
                    wanted_titles = [title]

                sheets_out = []
                for title in wanted_titles:
                    ws = wb[title]
                    raw_rows = list(ws.iter_rows(values_only=True))
                    # Coerce dates to ISO strings, leave numbers/bools as-is.
                    normalized_rows = []
                    for row in raw_rows:
                        normalized_rows.append([self._coerce_xlsx_cell(c) for c in row])
                    trimmed, truncated = self._truncate_rows(normalized_rows, max_rows)
                    columns = []
                    if trimmed:
                        columns = [str(c) if c is not None else "" for c in trimmed[0]]
                    sheets_out.append({
                        "name": title,
                        "row_count": len(normalized_rows),
                        "columns": columns,
                        "rows": trimmed,
                        "truncated": truncated,
                    })
                return {"path": normalized, "sheets": sheets_out}
            finally:
                try:
                    wb.close()
                except Exception:
                    pass
        except Exception as exc:
            return {"error": f"{type(exc).__name__}: {exc}", "path": normalized}

    @staticmethod
    def _coerce_xlsx_cell(cell):
        if cell is None:
            return None
        # openpyxl returns datetime / int / float / bool / str natively.
        import datetime as _dt
        if isinstance(cell, (_dt.datetime, _dt.date, _dt.time)):
            return cell.isoformat()
        return cell

    def read_docx(self, path: str) -> dict:
        """Read a DOCX file into paragraphs + headings + tables."""
        try:
            normalized = self._prepare_struct_read(path)
        except Exception as exc:
            return {"error": str(exc), "path": str(path)}

        try:
            import docx  # type: ignore
            doc = docx.Document(normalized)
            title = ""
            try:
                cp = doc.core_properties
                title = (cp.title or "").strip() if cp else ""
            except Exception:
                title = ""

            paragraphs = []
            headings = []
            for p in doc.paragraphs:
                text = (p.text or "").strip()
                style = (p.style.name if p.style else "") or ""
                paragraphs.append({"text": p.text or "", "style": style})
                if style.startswith("Heading") and text:
                    try:
                        level = int(style.split()[-1])
                    except (ValueError, IndexError):
                        level = 1
                    headings.append({"level": max(1, min(9, level)), "text": text})

            tables = []
            for tbl in doc.tables:
                rows = []
                for row in tbl.rows:
                    rows.append([cell.text for cell in row.cells])
                tables.append({"rows": rows})

            section_count = 0
            try:
                section_count = len(doc.sections)
            except Exception:
                section_count = 0

            return {
                "path": normalized,
                "title": title,
                "paragraphs": paragraphs,
                "headings": headings,
                "tables": tables,
                "section_count": int(section_count),
            }
        except Exception as exc:
            return {"error": f"{type(exc).__name__}: {exc}", "path": normalized}

    def read_pptx(self, path: str) -> dict:
        """Read a PPTX file into per-slide structured data."""
        try:
            normalized = self._prepare_struct_read(path)
        except Exception as exc:
            return {"error": str(exc), "path": str(path)}

        try:
            from pptx import Presentation  # type: ignore
            prs = Presentation(normalized)
            slides_out = []
            for idx, slide in enumerate(prs.slides, start=1):
                title_text = ""
                body_parts: list[str] = []
                tables_out: list[dict] = []
                notes_text = ""

                title_shape = None
                try:
                    title_shape = slide.shapes.title
                except Exception:
                    title_shape = None
                if title_shape is not None:
                    title_text = (title_shape.text or "").strip()

                for shape in slide.shapes:
                    if shape is title_shape:
                        continue
                    if getattr(shape, "has_table", False):
                        tbl = shape.table
                        rows = []
                        for row in tbl.rows:
                            rows.append([cell.text for cell in row.cells])
                        tables_out.append({"rows": rows})
                        continue
                    if getattr(shape, "has_text_frame", False):
                        txt = (shape.text or "").strip()
                        if txt:
                            body_parts.append(txt)

                try:
                    if slide.has_notes_slide:
                        notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
                except Exception:
                    notes_text = ""

                slides_out.append({
                    "index": idx,
                    "title": title_text,
                    "body_text": "\n".join(body_parts),
                    "tables": tables_out,
                    "notes": notes_text,
                })

            return {
                "path": normalized,
                "slide_count": len(slides_out),
                "slides": slides_out,
            }
        except Exception as exc:
            return {"error": f"{type(exc).__name__}: {exc}", "path": normalized}

    def list_local_directory(self, path: str):
        normalized = self._normalize_path(path)
        if os.path.isfile(normalized):
            normalized = os.path.dirname(normalized)
        if not os.path.exists(normalized):
            raise FileNotFoundError(normalized)
        if not os.path.isdir(normalized):
            raise NotADirectoryError(normalized)

        entries = []
        for name in sorted(os.listdir(normalized), key=str.lower):
            full_path = os.path.join(normalized, name)
            entries.append({
                "name": name,
                "path": full_path,
                "is_dir": os.path.isdir(full_path),
                "size": os.path.getsize(full_path) if os.path.isfile(full_path) else 0,
            })

        self.remember_paths([normalized] + [item["path"] for item in entries], source="local")
        return entries

    def write_user_file(
        self,
        path: str,
        content: str,
        *,
        overwrite: bool = True,
        mkdir: bool = True,
        encoding: str = "utf-8",
    ):
        """写入文本到本地文件。

        - 支持 `~`、`~/Desktop/foo.md` 等用户路径展开。
        - 默认允许创建父目录、覆盖已有文件。
        """
        if path is None or str(path).strip() == "":
            raise ValueError("path is required")
        if content is None:
            content = ""
        if not isinstance(content, str):
            content = str(content)

        # 不能直接用 _normalize_path，因为它会 resolve 不存在的父目录在某些情况下报错；
        # 用 expanduser + absolute 更稳。
        raw = pathlib.Path(self._coerce_path_value(path)).expanduser()
        if not raw.is_absolute():
            raw = raw.resolve()
        else:
            raw = pathlib.Path(os.path.normpath(str(raw)))

        if raw.exists() and raw.is_dir():
            raise IsADirectoryError(str(raw))
        if raw.exists() and not overwrite:
            raise FileExistsError(str(raw))

        parent = raw.parent
        if not parent.exists():
            if not mkdir:
                raise FileNotFoundError(str(parent))
            parent.mkdir(parents=True, exist_ok=True)

        if len(content.encode(encoding, errors="replace")) > self.MAX_TEXT_FILE_SIZE:
            raise ValueError(
                f"Content too large to write (> {self.MAX_TEXT_FILE_SIZE} bytes)"
            )

        with open(raw, "w", encoding=encoding, newline="") as f:
            f.write(content)

        normalized = str(raw.resolve())
        self.remember_paths([normalized], source="local")
        return {
            "path": normalized,
            "name": os.path.basename(normalized),
            "size": os.path.getsize(normalized),
            "encoding": encoding,
            "overwritten": False,  # 不强求精确判断
        }

    def build_file_tree_entries(
        self,
        base_dir: str,
        files: list[dict[str, Any]],
        *,
        encoding: str = "utf-8",
    ) -> dict[str, Any]:
        if not str(base_dir or "").strip():
            raise ValueError("base_dir is required")
        if not isinstance(files, list) or not files:
            raise ValueError("files must be a non-empty list")

        base = pathlib.Path(self._coerce_path_value(base_dir)).expanduser()
        base = base.resolve() if not base.is_absolute() else pathlib.Path(os.path.normpath(str(base)))
        base_norm = os.path.normcase(os.path.abspath(str(base)))
        entries = []
        total_bytes = 0
        seen_paths = set()

        for index, item in enumerate(files, start=1):
            if not isinstance(item, dict):
                raise ValueError(f"files[{index}] must be an object")
            relative_path = str(item.get("relative_path") or item.get("path") or "").strip()
            if not relative_path:
                raise ValueError(f"files[{index}].relative_path is required")
            rel = pathlib.PureWindowsPath(relative_path) if "\\" in relative_path else pathlib.PurePosixPath(relative_path)
            if rel.is_absolute() or ".." in rel.parts:
                raise ValueError(f"Unsafe relative path: {relative_path}")
            content = item.get("content", "")
            if content is None:
                content = ""
            if not isinstance(content, str):
                content = str(content)
            target = base.joinpath(*rel.parts)
            target_norm = os.path.normcase(os.path.abspath(str(target)))
            if os.path.commonpath([base_norm, target_norm]) != base_norm:
                raise ValueError(f"Path escapes base_dir: {relative_path}")
            if target_norm in seen_paths:
                raise ValueError(f"Duplicate file target: {relative_path}")
            seen_paths.add(target_norm)
            size = len(content.encode(encoding, errors="replace"))
            if size > self.MAX_TEXT_FILE_SIZE:
                raise ValueError(f"Content too large for {relative_path} (> {self.MAX_TEXT_FILE_SIZE} bytes)")
            total_bytes += size
            entries.append({
                "relative_path": str(pathlib.PurePosixPath(*rel.parts)),
                "path": str(target),
                "content": content,
                "size": size,
                "encoding": encoding,
            })

        return {
            "base_dir": str(base),
            "file_count": len(entries),
            "total_bytes": total_bytes,
            "files": entries,
        }

    def describe_local_path(self, path: str):
        normalized = self._normalize_path(path)
        if not os.path.exists(normalized):
            raise FileNotFoundError(normalized)

        info = {
            "path": normalized,
            "name": os.path.basename(normalized) or normalized,
            "exists": True,
            "is_dir": os.path.isdir(normalized),
            "is_file": os.path.isfile(normalized),
            "size": os.path.getsize(normalized) if os.path.isfile(normalized) else 0,
        }
        self.remember_paths([normalized], source="local")
        return info

    def get_selected_source_files(self):
        """Return local source files for the current Wwise selection.

        Strategy:
        1. Collect selected objects plus all descendants using IDs.
        2. Use `ak.wwise.core.object.get` with `select children` to retrieve child AudioFileSource objects.
        3. Read `originalFilePath` from `options.return`.
        """
        all_objects, warnings = self._get_selected_and_descendant_objects()
        if not all_objects:
            return []
        return self._get_source_files_from_sound_objects(all_objects, warnings)

    # --- Sound-object fallback for source file resolution ---

    _SOUND_LIKE_TYPES = {"sound", "soundsfx", "soundvoice", "musictrack"}

    def _get_source_files_from_sound_objects(self, all_objects, warnings):
        """Resolve source files by traversing to Sound-level objects first.

        Flow:
        1. Get selected objects and descendants by ID.
        2. Identify Sound-level objects from that set.
        3. Query only those Sound objects with `select children`.
        4. Read `originalFilePath` from AudioFileSource children.
        """
        candidate_ids = []
        id_to_obj = {}
        for item in all_objects:
            if not isinstance(item, dict):
                continue
            obj_id = (item.get("id") or "").strip()
            if obj_id and obj_id not in id_to_obj:
                candidate_ids.append(obj_id)
                id_to_obj[obj_id] = item

        if not candidate_ids:
            return []

        parent_info_items, parent_info_warnings = self._query_objects_by_ids_resilient(
            candidate_ids,
            ["id", "name", "type", "path", "activeSource", "childrenCount"],
        )
        warnings.extend(parent_info_warnings)
        for item in parent_info_items:
            if not isinstance(item, dict):
                continue
            item_id = (item.get("id") or "").strip()
            if item_id:
                id_to_obj[item_id] = {**id_to_obj.get(item_id, {}), **item}

        sound_objects = []
        for item in id_to_obj.values():
            if not isinstance(item, dict):
                continue
            item_id = (item.get("id") or "").strip()
            if not item_id:
                continue
            item_type = (item.get("type") or "").strip().casefold()
            active_source_id = self._extract_reference_id(item.get("activeSource"))
            if item_type in self._SOUND_LIKE_TYPES or active_source_id:
                sound_objects.append(item)

        sound_ids = []
        seen_sound_ids = set()
        for item in sound_objects:
            item_id = (item.get("id") or "").strip()
            if item_id and item_id not in seen_sound_ids:
                seen_sound_ids.add(item_id)
                sound_ids.append(item_id)

        if not sound_ids:
            warnings.append("未在选中对象及其后代中识别到 Sound 层级对象。")
            return []

        # Query children of Sound objects only to get AudioFileSource items.
        result = self.waapi_client.call(
            "ak.wwise.core.object.get",
            args={
                "from": {"id": sound_ids},
                "transform": [{"select": ["children"]}],
            },
            options={"return": ["id", "name", "type", "path", "parent", "originalFilePath"]},
        )
        error_text = self._extract_waapi_error(result)
        if error_text:
            warnings.append(f"查询对象子项失败: {error_text}")
            return []

        items = result.get("return", []) if isinstance(result, dict) else []

        audio_source_children = []
        for obj in (items if isinstance(items, list) else []):
            if not isinstance(obj, dict):
                continue
            obj_type = (obj.get("type") or "").strip().casefold()
            has_file = bool((obj.get("originalFilePath") or "").strip())
            if obj_type == "audiofilesource" or (not obj_type and has_file):
                audio_source_children.append(obj)

        files = []
        seen_paths = set()
        for obj in audio_source_children:
            source_id = (obj.get("id") or "").strip()
            parent_id = self._extract_reference_id(obj.get("parent"))
            parent_item = id_to_obj.get(parent_id, {}) if parent_id else {}
            active_source_id = self._extract_reference_id(parent_item.get("activeSource"))
            if active_source_id and active_source_id != source_id:
                continue

            wav = self._resolve_project_source_path(obj.get("originalFilePath") or "")
            if not wav or not os.path.isfile(wav):
                continue
            if wav in seen_paths:
                continue
            seen_paths.add(wav)

            entry = self._make_project_source_entry(
                wav,
                object_name=(parent_item.get("name") or obj.get("name") or ""),
                object_id=(parent_item.get("id") or source_id or ""),
                object_path=(parent_item.get("path") or obj.get("path") or ""),
                object_type=(parent_item.get("type") or ""),
                source_object_type="AudioFileSource",
            )
            entry["sourceObjectId"] = source_id
            entry["sourceObjectName"] = (obj.get("name") or "").strip()
            entry["sourceObjectPath"] = (obj.get("path") or "").strip()
            entry["originalWavFilePath"] = entry["originalFilePath"]
            entry["filePath"] = entry["originalFilePath"]
            entry["filepath"] = entry["originalFilePath"]
            files.append(entry)
        return files

    def get_selected_source_filepaths(self):
        """Return local source audio file paths for current Wwise selection.

        This is a convenience wrapper around `get_selected_source_files()` when
        only the physical file paths are needed.
        """
        source_files = self.get_selected_source_files()
        paths = []
        seen = set()
        for item in source_files:
            if not isinstance(item, dict):
                continue
            path = (item.get("originalFilePath") or item.get("filePath") or item.get("path") or "").strip()
            if not path or path in seen:
                continue
            seen.add(path)
            paths.append(path)
        return paths

    def import_audio_files_to_selected_wwise(
        self,
        paths,
        *,
        object_type: str = "Sound SFX",
        import_operation: str = "useExisting",
        import_language: str = "SFX",
        originals_sub_folder: str = "",
        wait_for_files: bool = True,
        file_ready_timeout: float = 30.0,
        batch_size: int = 20,
        retry_on_copy_failure: bool = True,
    ):
        """Import local audio files as Sound objects under the current Wwise selection."""
        if isinstance(paths, (str, pathlib.Path, dict)):
            raw_paths = [paths]
        else:
            raw_paths = list(paths or [])
        if not raw_paths:
            raise ValueError("No audio files were provided for Wwise import.")

        selected_objects, warnings = self._get_selected_wwise_objects()
        if not selected_objects:
            raise RuntimeError("Wwise 中没有可用的选中层级，无法导入音频。")
        target = selected_objects[0]
        target_id = (target.get("id") or "").strip()
        if not target_id:
            raise RuntimeError("当前 Wwise 选中对象缺少 id，无法作为导入目标。")

        supported_extensions = {".wav", ".wave", ".aif", ".aiff", ".flac", ".ogg", ".mp3", ".wem"}
        imports = []
        files = []
        skipped = []
        seen_paths = set()
        for raw_path in raw_paths:
            try:
                normalized = self._normalize_path(raw_path)
            except Exception as exc:
                skipped.append({"path": str(raw_path), "reason": str(exc)})
                continue
            if normalized in seen_paths:
                continue
            seen_paths.add(normalized)
            suffix = os.path.splitext(normalized)[1].casefold()
            if suffix not in supported_extensions:
                skipped.append({"path": normalized, "reason": "不是支持的音频文件类型。"})
                continue
            if not os.path.isfile(normalized):
                skipped.append({"path": normalized, "reason": "文件不存在。"})
                continue
            item = {
                "audioFile": normalized,
                "objectPath": self._wwise_object_name_from_file(normalized),
                "objectType": object_type or "Sound SFX",
            }
            if originals_sub_folder:
                item["originalsSubFolder"] = str(originals_sub_folder).strip().strip("\\/")
            imports.append(item)
            files.append(normalized)

        if not imports:
            return {
                "ok": False,
                "target": target,
                "imported_count": 0,
                "requested_count": len(raw_paths),
                "skipped": skipped,
                "warnings": [*warnings, "没有可导入的有效音频文件。"],
            }

        if wait_for_files:
            ready_files, not_ready = self._wait_for_files_ready(files, timeout=file_ready_timeout)
            skipped.extend(not_ready)
            ready_set = set(ready_files)
            imports = [item for item in imports if item.get("audioFile") in ready_set]
            files = [path for path in files if path in ready_set]
            if not imports:
                return {
                    "ok": False,
                    "target": target,
                    "imported_count": 0,
                    "requested_count": len(raw_paths),
                    "source_files": [],
                    "imported_files": [],
                    "objects": [],
                    "log": [],
                    "skipped": skipped,
                    "warnings": [*warnings, "渲染文件尚未稳定或不可读，已取消导入。"],
                    "error": "渲染文件尚未稳定或不可读。",
                    "raw_result": {},
                }

        operation = import_operation if import_operation in {"createNew", "useExisting", "replaceExisting"} else "useExisting"
        default = {"importLocation": target_id}
        if import_language:
            default["importLanguage"] = str(import_language)
        results = []
        log_items = []
        imported_files = []
        imported_objects = []
        error_logs = []
        raw_errors = []
        for batch in self._chunked(imports, batch_size):
            result = self._call_wwise_audio_import(operation, default, batch)
            if self._result_has_copy_failure(result) and retry_on_copy_failure:
                retry_paths = [item.get("audioFile") for item in batch if item.get("audioFile")]
                self._wait_for_files_ready(retry_paths, timeout=min(max(float(file_ready_timeout or 0), 5.0), 15.0))
                result = self._call_wwise_audio_import(operation, default, batch)
            results.append(result)
            error_text = self._extract_waapi_error(result)
            if error_text:
                raw_errors.append(error_text)
            batch_log = result.get("log", []) if isinstance(result, dict) and isinstance(result.get("log"), list) else []
            log_items.extend(batch_log)
            if isinstance(result, dict):
                if isinstance(result.get("files"), list):
                    imported_files.extend(result.get("files") or [])
                if isinstance(result.get("objects"), list):
                    imported_objects.extend(result.get("objects") or [])
        error_text = "; ".join(item for item in raw_errors if item)
        error_logs = []
        for entry in log_items:
            if not isinstance(entry, dict):
                continue
            severity = str(entry.get("severity") or "").casefold()
            message = str(entry.get("message") or "").strip()
            if severity in {"error", "fatal"} or "error" in severity:
                error_logs.append(message or str(entry))
        ok = not error_text and not error_logs
        self.remember_paths(files, source="wwise-import")
        return {
            "ok": ok,
            "target": {
                "id": target_id,
                "name": target.get("name", ""),
                "type": target.get("type", ""),
                "path": target.get("path", ""),
            },
            "imported_count": len(imported_files) or len(imported_objects),
            "requested_count": len(raw_paths),
            "source_files": files,
            "imported_files": imported_files,
            "objects": imported_objects,
            "log": log_items if isinstance(log_items, list) else [],
            "skipped": skipped,
            "warnings": warnings,
            "error": error_text or "; ".join(error_logs),
            "raw_result": results[0] if len(results) == 1 and isinstance(results[0], dict) else {"batches": results},
        }

    def _call_wwise_audio_import(self, operation: str, default: dict[str, Any], imports: list[dict[str, Any]]) -> dict[str, Any]:
        try:
            result = self.waapi_client.call(
                "ak.wwise.core.audio.import",
                args={
                    "importOperation": operation,
                    "default": default,
                    "imports": imports,
                },
                options={"return": ["id", "name", "type", "path", "originalFilePath"]},
            )
        except Exception as exc:
            return self._waapi_exception_to_result(exc)
        return result if isinstance(result, dict) else {"error": str(result)}

    @staticmethod
    def _waapi_exception_to_result(exc: Exception) -> dict[str, Any]:
        payload = {}
        kwargs = getattr(exc, "kwargs", None)
        if isinstance(kwargs, dict):
            payload = kwargs
        details = payload.get("details") if isinstance(payload, dict) else None
        log = details.get("log") if isinstance(details, dict) and isinstance(details.get("log"), list) else []
        message = payload.get("message") if isinstance(payload, dict) else ""
        return {
            "error": str(message or exc),
            "log": log,
            "files": [],
            "objects": [],
            "exception": type(exc).__name__,
            "details": details if isinstance(details, dict) else {},
        }

    @staticmethod
    def _result_has_copy_failure(result: dict[str, Any]) -> bool:
        if not isinstance(result, dict):
            return False
        log_items = result.get("log") if isinstance(result.get("log"), list) else []
        return any("copy file to originals folder failed" in str(item.get("message") or "").casefold() for item in log_items if isinstance(item, dict))

    def analyze_selected_source_files_loudness(self, limit: int | None = None, source_files=None):
        if source_files is None:
            source_files = self.get_selected_source_files()
        if isinstance(limit, int) and limit > 0:
            source_files = source_files[:limit]

        results = []
        warnings = []
        total = len(source_files)
        try:
            for index, item in enumerate(source_files, start=1):
                if not isinstance(item, dict):
                    continue
                source_path = (item.get("originalFilePath") or item.get("path") or "").strip()
                object_name = (item.get("objectName") or item.get("name") or source_path or "unknown").strip()
                self._report_analysis_progress(index, total, object_name)
                if not source_path or not os.path.isfile(source_path):
                    warnings.append(f"{object_name}: 本地文件不存在，已跳过。")
                    continue
                try:
                    analysis = self.analyze_audio_file(source_path)
                except Exception as exc:
                    warnings.append(f"{object_name}: 响度分析失败: {exc}")
                    continue
                row = dict(analysis)
                row.update({
                    "objectName": object_name,
                    "objectId": item.get("objectId", ""),
                    "objectPath": item.get("objectPath", ""),
                    "originalFilePath": source_path,
                    "file": os.path.basename(source_path),
                    "file_name": os.path.basename(source_path),
                    "duration": analysis.get("duration_seconds"),
                    "integrated_loudness": analysis.get("integrated_loudness_lufs"),
                    "loudness_range": analysis.get("loudness_range_lu"),
                    "analysis": analysis,
                })
                results.append(row)
        finally:
            self._finish_analysis_progress()

        return AnalysisReport({
            "count": len(results),
            "results": results,
            "warnings": warnings,
        })

    def analyze_directory_loudness(
        self,
        path: str,
        *,
        recursive: bool = True,
        extensions: list[str] | None = None,
        limit: int | None = None,
        top_n_frequencies: int = DEFAULT_DOMINANT_FREQUENCIES,
    ):
        root = self._normalize_path(path)
        if not os.path.exists(root):
            raise FileNotFoundError(root)
        normalized_exts = {str(ext).lower() if str(ext).startswith(".") else f".{str(ext).lower()}" for ext in (extensions or [".wav"])}
        if os.path.isfile(root):
            if os.path.splitext(root)[1].lower() not in normalized_exts:
                raise ValueError(f"File extension is not in allowed extensions: {root}")
            files = [root]
            root_path = os.path.dirname(root)
        elif os.path.isdir(root):
            root_path = root
            files = []
            if recursive:
                for current_root, _dirs, names in os.walk(root):
                    for name in names:
                        full = os.path.join(current_root, name)
                        if os.path.splitext(name)[1].lower() in normalized_exts:
                            files.append(full)
            else:
                for name in os.listdir(root):
                    full = os.path.join(root, name)
                    if os.path.isfile(full) and os.path.splitext(name)[1].lower() in normalized_exts:
                        files.append(full)
            files = sorted(files, key=str.lower)
        else:
            raise ValueError(f"Unsupported local path: {root}")

        if isinstance(limit, int) and limit > 0:
            files = files[:limit]

        results = []
        warnings = []
        total = len(files)
        try:
            for index, file_path in enumerate(files, start=1):
                self._report_analysis_progress(index, total, os.path.basename(file_path))
                try:
                    analysis = self.analyze_audio_file(file_path, top_n_frequencies=top_n_frequencies)
                except Exception as exc:
                    warnings.append(f"{file_path}: 响度分析失败: {exc}")
                    continue
                row = dict(analysis)
                row.update({
                    "file": os.path.basename(file_path),
                    "file_name": os.path.basename(file_path),
                    "relative_path": os.path.relpath(file_path, root_path),
                    "path": analysis.get("path", file_path),
                    "duration": analysis.get("duration_seconds"),
                    "integrated_loudness": analysis.get("integrated_loudness_lufs"),
                    "loudness_range": analysis.get("loudness_range_lu"),
                })
                results.append(row)
        finally:
            self._finish_analysis_progress()

        finite_lufs = [item.get("integrated_loudness_lufs") for item in results if self._is_finite_number(item.get("integrated_loudness_lufs"))]
        finite_peaks = [item.get("true_peak_dbfs") for item in results if self._is_finite_number(item.get("true_peak_dbfs"))]
        total_duration = sum(float(item.get("duration_seconds") or 0) for item in results)
        summary = {
            "analyzed_count": len(results),
            "total_duration_seconds": round(total_duration, 4),
            "average_lufs_i": self._round_or_none(sum(finite_lufs) / len(finite_lufs)) if finite_lufs else None,
            "max_lufs_i": self._round_or_none(max(finite_lufs)) if finite_lufs else None,
            "min_lufs_i": self._round_or_none(min(finite_lufs)) if finite_lufs else None,
            "max_true_peak_dbfs": self._round_or_none(max(finite_peaks)) if finite_peaks else None,
        }
        if total == 0:
            warnings.append("未找到匹配的音频文件。")

        self.remember_paths([root_path] + files, source="local")
        return AnalysisReport({
            "root_path": root_path,
            "file_count": total,
            "analyzed_count": len(results),
            "count": len(results),
            "results": results,
            "warnings": warnings,
            "summary": summary,
        })

    def get_project_source_files(self, object_ids: list = None, object_type: str = "Sound"):
        """Return locally resolvable project source files.

        Broad project scans prefer source-control listing because some Wwise versions do
        not expose `originalFilePath` on `ak.wwise.core.object.get`.
        """
        if not object_ids:
            source_entries = self._list_source_control_project_entries()
            if source_entries:
                return source_entries

        if object_ids:
            returned, _warnings, has_original_path = self._query_objects_with_optional_source_path(object_ids, ["id", "name", "path"])
        else:
            result = self.waapi_client.call(
                "ak.wwise.core.object.get",
                args={"from": {"ofType": [object_type]}},
                options={"return": ["id", "name", "path"]},
            )
            if not isinstance(result, dict):
                return []

            returned = result.get("return", [])
            if not isinstance(returned, list):
                return []
            has_original_path = False

        remembered = []
        if has_original_path:
            for obj in returned:
                if not isinstance(obj, dict):
                    continue
                wav = (obj.get("originalFilePath") or "").strip()
                if wav and os.path.isfile(wav):
                    remembered.append(
                        self._make_project_source_entry(
                            wav,
                            object_name=(obj.get("name") or ""),
                            object_id=(obj.get("id") or ""),
                            object_path=(obj.get("path") or ""),
                            object_type=(obj.get("type") or object_type),
                            source_object_type=(obj.get("type") or object_type),
                        )
                    )

        if remembered or object_ids:
            if remembered:
                return remembered
            fallback_results, _match_warnings = self._match_objects_to_source_entries(returned, self._list_source_control_project_entries())
            return fallback_results

        return self._list_source_control_project_entries()

    def analyze_audio_file(self, path: str, top_n_frequencies: int = DEFAULT_DOMINANT_FREQUENCIES):
        normalized = self._normalize_path(path)
        if not os.path.exists(normalized):
            raise FileNotFoundError(normalized)
        if os.path.isdir(normalized):
            raise IsADirectoryError(normalized)

        self.remember_paths([normalized], source="local")

        np, librosa, pyln, sf = self._import_audio_analysis_dependencies()
        audio, sample_rate, source_info = self._load_audio_samples(normalized, np, librosa, sf)
        if audio.size == 0:
            raise ValueError("Audio file is empty or could not be decoded.")

        frame_count = int(audio.shape[0])
        channels = int(audio.shape[1]) if audio.ndim > 1 else 1
        mono_audio = audio[:, 0] if channels == 1 else np.mean(audio, axis=1)
        duration_seconds = frame_count / float(sample_rate) if sample_rate else 0.0
        peak = float(np.max(np.abs(audio))) if audio.size else 0.0
        rms = float(np.sqrt(np.mean(np.square(audio)))) if audio.size else 0.0
        peak_dbfs = self._linear_to_dbfs(peak)
        rms_dbfs = self._linear_to_dbfs(rms)

        # Real ITU-R BS.1770 inter-sample (true) peak via 4x oversampling.
        # Falls back to the sample peak if oversampling is unavailable.
        true_peak_dbfs = self._compute_true_peak_dbfs(audio, sample_rate, np, librosa)
        if true_peak_dbfs is None:
            true_peak_dbfs = peak_dbfs

        loudness_metrics = self._analyze_loudness(audio, sample_rate, pyln, np)
        integrated_lufs = loudness_metrics["integrated_lufs"]
        momentary_loudness_max_lufs = loudness_metrics["momentary_loudness_max_lufs"]
        loudness_range_lu = loudness_metrics["loudness_range_lu"]

        integrated_lufs = self._round_or_none(integrated_lufs)
        momentary_loudness_max_lufs = self._round_or_none(momentary_loudness_max_lufs)
        loudness_range_lu = self._round_or_none(loudness_range_lu)

        spectral_summary = self._analyze_spectrum(
            mono_audio,
            sample_rate,
            top_n_frequencies=top_n_frequencies,
            np=np,
            librosa=librosa,
        )

        return {
            "path": normalized,
            "analysis_backend": source_info["backend"],
            "format": source_info.get("format", "unknown"),
            "subtype": source_info.get("subtype", "unknown"),
            "channels": channels,
            "sample_rate": int(sample_rate),
            "frame_count": frame_count,
            "duration_seconds": round(duration_seconds, 4),
            "duration": round(duration_seconds, 4),
            "peak_normalized": round(peak, 6),
            "rms_normalized": round(rms, 6),
            "peak_dbfs": peak_dbfs,
            "rms_dbfs": rms_dbfs,
            "true_peak_dbfs": true_peak_dbfs,
            "true_peak": true_peak_dbfs,
            "integrated_loudness_lufs": integrated_lufs,
            "lufs_i": integrated_lufs,
            "integrated_loudness": integrated_lufs,
            "momentary_loudness_max_lufs": momentary_loudness_max_lufs,
            "lufs_m_max": momentary_loudness_max_lufs,
            "loudness_range_lu": loudness_range_lu,
            "loudness_range": loudness_range_lu,
            "analysis_warnings": loudness_metrics["warnings"],
            **spectral_summary,
        }

    def analyze_wav_file(self, path: str):
        normalized = self._normalize_path(path)
        if not normalized.lower().endswith(".wav"):
            raise ValueError("Only WAV analysis is supported in this version.")
        result = self.analyze_audio_file(normalized)
        result["analysis_tool"] = "analyze_wav_file"
        return result

    @staticmethod
    def _import_audio_analysis_dependencies():
        try:
            import numpy as np
            import librosa
            import pyloudnorm as pyln
            import soundfile as sf
        except ImportError as exc:
            raise RuntimeError(
                "Audio analysis dependencies are unavailable. Please install pyloudnorm, librosa and soundfile."
            ) from exc
        return np, librosa, pyln, sf

    @staticmethod
    def _load_audio_samples(path: str, np, librosa, sf):
        try:
            info = sf.info(path)
            audio, sample_rate = sf.read(path, dtype="float32", always_2d=True)
            return audio, sample_rate, {
                "backend": "soundfile",
                "format": getattr(info, "format", "unknown"),
                "subtype": getattr(info, "subtype", "unknown"),
            }
        except Exception:
            loaded, sample_rate = librosa.load(path, sr=None, mono=False)
            if loaded.ndim == 1:
                audio = loaded.reshape(-1, 1)
            else:
                audio = loaded.T
            return np.asarray(audio, dtype=np.float32), sample_rate, {
                "backend": "librosa",
                "format": os.path.splitext(path)[1].lower().lstrip(".") or "unknown",
                "subtype": "decoded",
            }

    @staticmethod
    def _is_finite_number(value) -> bool:
        try:
            return math.isfinite(float(value))
        except (TypeError, ValueError):
            return False

    @staticmethod
    def _coerce_meter_input(audio, np):
        meter_input = np.asarray(audio, dtype=np.float64)
        if meter_input.ndim == 2 and meter_input.shape[1] == 1:
            return meter_input[:, 0]
        return meter_input

    @staticmethod
    def _channel_count_of(meter_input) -> int:
        return int(meter_input.shape[1]) if getattr(meter_input, "ndim", 1) == 2 else 1

    def _downmix_to_measurable(self, meter_input, np):
        """pyloudnorm (ITU-R BS.1770) supports at most 5 channels. Audio with
        more channels (e.g. 5.1 = 6ch, 7.1 = 8ch — very common in game audio)
        makes ``meter.integrated_loudness`` raise ``ValueError`` via
        ``util.valid_audio``. Rather than swallow that and report a null/wrong
        cause, downmix to mono so a value can still be produced.

        Returns ``(array, original_channels, downmixed)``.
        """
        channels = self._channel_count_of(meter_input)
        if channels <= 5:
            return meter_input, channels, False
        mono = np.mean(np.asarray(meter_input, dtype=np.float64), axis=1)
        return mono, channels, True

    def _compute_true_peak_dbfs(self, audio, sample_rate, np, librosa, oversample: int = 4):
        """ITU-R BS.1770 inter-sample (true) peak via oversampling.

        ``audio`` is the raw ``(frames, channels)`` float array. Each channel is
        upsampled by ``oversample`` (default 4x) and the max absolute inter-sample
        value is taken. Returns dBFS (or ``None`` if it cannot be computed — the
        caller falls back to the sample peak).
        """
        try:
            if not sample_rate or sample_rate <= 0:
                return None
            data = np.asarray(audio, dtype=np.float32)
            if data.size == 0:
                return None
            if data.ndim == 1:
                data = data.reshape(-1, 1)
            target_sr = int(sample_rate) * int(oversample)
            peak = 0.0
            for ch in range(data.shape[1]):
                column = data[:, ch]
                try:
                    upsampled = librosa.resample(column, orig_sr=sample_rate, target_sr=target_sr)
                except Exception:
                    upsampled = librosa.resample(column, orig_sr=sample_rate, target_sr=target_sr, res_type="fft")
                ch_peak = float(np.max(np.abs(upsampled))) if upsampled.size else 0.0
                if ch_peak > peak:
                    peak = ch_peak
            return self._linear_to_dbfs(peak)
        except Exception:
            return None


    def _extract_blockwise_loudness_values(self, meter) -> list[float]:
        blockwise_loudness = getattr(meter, "blockwise_loudness", None) or []
        if callable(blockwise_loudness):
            return []
        return [float(value) for value in blockwise_loudness if self._is_finite_number(value)]

    def _measure_loudness_with_fallback(self, meter_input, sample_rate: int, pyln, np):
        """Measure LUFS-I with fallback to LUFS-M Max for short audio (<400ms).

        Returns:
            (lufs_value, momentary_max, warnings) where lufs_value is the
            integrated loudness (or LUFS-M Max approximation for short audio),
            momentary_max is the LUFS-M Max value, and warnings is a list of
            warning strings.
        """
        warnings: list[str] = []

        # --- 0. Channel guard: pyloudnorm/BS.1770 supports <=5 channels. ---
        # Downmix higher channel counts (5.1/7.1 beds are common in game audio)
        # to mono so we still return a value instead of a misleading null.
        meter_input, original_channels, downmixed = self._downmix_to_measurable(meter_input, np)
        if downmixed:
            warnings.append(
                f"音频为 {original_channels} 声道，超过 LUFS 测量上限（5 声道），"
                "已降混为单声道后测量，结果为近似值。"
            )

        duration_seconds = float(meter_input.shape[0]) / float(sample_rate) if sample_rate else 0.0
        integrated_lufs = None
        momentary_max = None
        used_fallback = False

        # --- 1. Standard integrated loudness ---
        meter = pyln.Meter(sample_rate)
        try:
            integrated_lufs = float(meter.integrated_loudness(meter_input))
            finite_blocks = self._extract_blockwise_loudness_values(meter)
            if finite_blocks:
                momentary_max = max(finite_blocks)
        except ValueError:
            pass  # handled below via fallback

        # --- 2. Short-audio fallback: pad to 400ms, take LUFS-M Max ---
        # Zero-padding to 400ms dilutes the block energy by sample_count/min_samples,
        # biasing the LUFS-M Max low by 10*log10(sample_count/min_samples) dB. We
        # add the inverse correction so the value reflects the real signal level.
        if momentary_max is None and duration_seconds > 0:
            min_samples = max(1, int(round(sample_rate * 0.4)))
            sample_count = meter_input.shape[0]
            if sample_count < min_samples:
                if meter_input.ndim == 1:
                    padded = np.pad(meter_input, (0, min_samples - sample_count))
                else:
                    padded = np.pad(meter_input, ((0, min_samples - sample_count), (0, 0)))
                try:
                    pad_meter = pyln.Meter(sample_rate)
                    _ = pad_meter.integrated_loudness(padded)
                    finite_blocks = self._extract_blockwise_loudness_values(pad_meter)
                    if finite_blocks:
                        pad_correction_db = 10.0 * math.log10(min_samples / sample_count)
                        momentary_max = max(finite_blocks) + pad_correction_db
                except ValueError:
                    pass

        # --- 3. If standard LUFS-I failed, use LUFS-M Max as fallback ---
        if not self._is_finite_number(integrated_lufs):
            if self._is_finite_number(momentary_max):
                integrated_lufs = momentary_max
                used_fallback = True
                warnings.append(
                    "音频短于 400ms，LUFS-I 使用 LUFS-M Max（补零并已做能量补偿）替代。"
                )
            else:
                if duration_seconds < 0.4:
                    warnings.append("音频短于 400ms，无法计算 LUFS-I。")
                else:
                    warnings.append("无法计算标准 LUFS-I。")

        if self._is_finite_number(momentary_max) and used_fallback:
            warnings.append("音频短于 400ms，LUFS-M Max 为补零到 400ms 并做能量补偿后的近似值。")

        return integrated_lufs, momentary_max, warnings

    def _analyze_loudness(self, audio, sample_rate: int, pyln, np) -> dict:
        meter_input = self._coerce_meter_input(audio, np)
        duration_seconds = float(audio.shape[0]) / float(sample_rate) if sample_rate else 0.0

        integrated_lufs, momentary_loudness_max_lufs, warnings = \
            self._measure_loudness_with_fallback(meter_input, sample_rate, pyln, np)

        loudness_range_lu = None
        # LRA also goes through pyloudnorm's <=5ch validation; downmix the same
        # way so multichannel beds still yield a value (the channel warning is
        # already emitted by _measure_loudness_with_fallback).
        lra_input, _orig_ch, _dm = self._downmix_to_measurable(meter_input, np)
        lra_meter = pyln.Meter(sample_rate)
        try:
            loudness_range_lu = float(lra_meter.loudness_range(lra_input))
        except ValueError:
            if duration_seconds < 3.0:
                warnings.append("音频短于约 3 秒，LRA 不具统计意义，已返回空值。")
            else:
                warnings.append("无法计算 LRA。")

        return {
            "integrated_lufs": self._round_or_none(integrated_lufs),
            "momentary_loudness_max_lufs": self._round_or_none(momentary_loudness_max_lufs),
            "loudness_range_lu": self._round_or_none(loudness_range_lu),
            "warnings": warnings,
        }

    def _round_or_none(self, value, digits: int = 2):
        if not self._is_finite_number(value):
            return None
        return round(float(value), digits)

    def _linear_to_dbfs(self, value, digits: int = 2):
        if not self._is_finite_number(value) or float(value) <= 0:
            return None
        return round(20.0 * math.log10(float(value)), digits)

    @staticmethod
    def _analysis_fft_size(sample_count: int) -> int:
        if sample_count <= 0:
            raise ValueError("Audio file contains no samples.")
        upper_bound = min(4096, sample_count)
        if upper_bound < 32:
            return upper_bound
        return 1 << int(math.floor(math.log2(upper_bound)))

    def _analyze_spectrum(self, mono_audio, sample_rate: int, top_n_frequencies: int, np, librosa):
        n_fft = self._analysis_fft_size(len(mono_audio))
        if n_fft < 16:
            raise ValueError("Audio file is too short for spectral analysis.")

        hop_length = max(8, n_fft // 4)
        centroid = float(np.mean(librosa.feature.spectral_centroid(y=mono_audio, sr=sample_rate, n_fft=n_fft, hop_length=hop_length)))
        bandwidth = float(np.mean(librosa.feature.spectral_bandwidth(y=mono_audio, sr=sample_rate, n_fft=n_fft, hop_length=hop_length)))
        rolloff = float(np.mean(librosa.feature.spectral_rolloff(y=mono_audio, sr=sample_rate, n_fft=n_fft, hop_length=hop_length, roll_percent=0.85)))
        flatness = float(np.mean(librosa.feature.spectral_flatness(y=mono_audio, n_fft=n_fft, hop_length=hop_length)))
        zcr = float(np.mean(librosa.feature.zero_crossing_rate(y=mono_audio, hop_length=hop_length)))

        stft_magnitude = np.abs(librosa.stft(mono_audio, n_fft=n_fft, hop_length=hop_length))
        averaged_spectrum = np.mean(stft_magnitude, axis=1)
        freqs = librosa.fft_frequencies(sr=sample_rate, n_fft=n_fft)
        dominant_indexes = np.argsort(averaged_spectrum)[-max(1, int(top_n_frequencies)):][::-1]
        dominant_frequencies = [round(float(freqs[idx]), 2) for idx in dominant_indexes if averaged_spectrum[idx] > 0]

        return {
            "spectral_centroid_hz": round(centroid, 2),
            "spectral_bandwidth_hz": round(bandwidth, 2),
            "spectral_rolloff_hz": round(rolloff, 2),
            "spectral_flatness": round(flatness, 6),
            "zero_crossing_rate": round(zcr, 6),
            "dominant_frequencies_hz": dominant_frequencies,
        }

    def _safe_db_value(self, value) -> float:
        try:
            as_float = float(value)
        except (TypeError, ValueError):
            return 0.0
        return as_float if math.isfinite(as_float) else 0.0

    @staticmethod
    def _extract_waapi_error(result) -> str:
        if isinstance(result, dict):
            error_text = result.get("error")
            if isinstance(error_text, str):
                return error_text.strip()
        if isinstance(result, str):
            return result.strip()
        return ""

    @staticmethod
    def _is_unknown_object_error(error_text: str) -> bool:
        normalized = (error_text or "").lower()
        return (
            "unknown_object" in normalized
            or "object is unknown" in normalized
            or "from id object is unknown" in normalized
        )

    @staticmethod
    def _is_unknown_property_error(error_text: str, property_name: str = "") -> bool:
        normalized = (error_text or "").lower()
        if "unknown property" not in normalized:
            return False
        property_token = (property_name or "").strip().lower()
        return not property_token or property_token in normalized

    def _query_objects_with_optional_source_path(self, object_ids: list[str], base_fields: list[str]):
        requested_fields = list(dict.fromkeys([*(base_fields or []), "originalFilePath"]))
        result = self.waapi_client.call(
            "ak.wwise.core.object.get",
            args={"from": {"id": object_ids}},
            options={"return": requested_fields},
        )
        error_text = self._extract_waapi_error(result)
        if not error_text:
            items = result.get("return", []) if isinstance(result, dict) else []
            return [item for item in items if isinstance(item, dict)], [], True

        if self._is_unknown_object_error(error_text):
            recovered_items, warnings = self._query_objects_by_ids_resilient(object_ids, requested_fields)
            return recovered_items, warnings, True

        if self._is_unknown_property_error(error_text, "originalfilepath"):
            basic_items, warnings = self._query_objects_by_ids_resilient(object_ids, base_fields)
            return basic_items, ["当前调用未返回 `originalFilePath`，已自动切换到兼容模式。", *warnings], False

        return [], [f"查询对象失败: {error_text}"], False

    def _query_objects_by_ids_resilient(self, object_ids: list[str], return_fields: list[str]):
        warnings = []
        if not object_ids:
            return [], warnings

        result = self.waapi_client.call(
            "ak.wwise.core.object.get",
            args={"from": {"id": object_ids}},
            options={"return": return_fields},
        )
        error_text = self._extract_waapi_error(result)
        if not error_text:
            items = result.get("return", []) if isinstance(result, dict) else []
            return [item for item in items if isinstance(item, dict)], warnings

        if not self._is_unknown_object_error(error_text):
            warnings.append(f"批量查询对象失败: {error_text}")
            return [], warnings

        warnings.append("检测到无效或陈旧的 Wwise 对象 ID，已自动跳过不可用对象。")
        recovered_items = []
        for object_id in object_ids:
            single_result = self.waapi_client.call(
                "ak.wwise.core.object.get",
                args={"from": {"id": [object_id]}},
                options={"return": return_fields},
            )
            single_error = self._extract_waapi_error(single_result)
            if single_error:
                if self._is_unknown_object_error(single_error):
                    warnings.append(f"{object_id}: 对象已失效或已不存在，已跳过。")
                else:
                    warnings.append(f"{object_id}: 查询失败: {single_error}")
                continue

            items = single_result.get("return", []) if isinstance(single_result, dict) else []
            if not items:
                warnings.append(f"{object_id}: 查询结果为空，已跳过。")
                continue

            first_item = items[0]
            if isinstance(first_item, dict):
                recovered_items.append(first_item)

        return recovered_items, warnings

    @staticmethod
    def _extract_reference_id(value):
        if isinstance(value, str):
            stripped = value.strip()
            return stripped or None
        if isinstance(value, dict):
            for key in ("id", "value", "guid"):
                candidate = value.get(key)
                if isinstance(candidate, str) and candidate.strip():
                    return candidate.strip()
            nested = value.get("object")
            if isinstance(nested, dict):
                nested_id = nested.get("id")
                if isinstance(nested_id, str) and nested_id.strip():
                    return nested_id.strip()
        if isinstance(value, list):
            for item in value:
                candidate = AgentToolbox._extract_reference_id(item)
                if candidate:
                    return candidate
        return None

    def _fetch_object_route_node(self, object_id: str):
        result = self.waapi_client.call(
            "ak.wwise.core.object.get",
            args={"from": {"id": [object_id]}},
            options={
                "return": [
                    "id",
                    "name",
                    "type",
                    "path",
                    "parent",
                    "@Volume",
                    "@BusVolume",
                    "@MakeUpGain",
                    "@OutputBus",
                    "OutputBus",
                ]
            },
        )
        error_text = self._extract_waapi_error(result)
        if error_text:
            if self._is_unknown_object_error(error_text):
                return {
                    "node": None,
                    "warning": f"{object_id}: 路由节点已失效或已不存在，停止继续向上追踪。",
                }
            return {
                "node": None,
                "warning": f"{object_id}: 查询路由节点失败: {error_text}",
            }
        if not isinstance(result, dict):
            return {"node": None, "warning": f"{object_id}: 路由节点查询返回格式无效。"}
        items = result.get("return", [])
        if not items:
            return {"node": None, "warning": f"{object_id}: 路由节点查询结果为空。"}
        return {"node": items[0], "warning": ""}

    def _estimate_ancestor_gain_chain(self, object_id: str):
        """Query ancestors of *object_id* and sum @Volume + @MakeUpGain.

        Also returns the effective OutputBus id (from the object itself or
        inherited from the nearest ancestor that defines one).
        """
        warnings: list[str] = []
        ancestors_list: list[dict] = []
        total_volume_db = 0.0
        total_makeup_db = 0.0
        effective_bus_id: str | None = None

        # --- query the object itself for its @OutputBus ---------
        obj_result = self.waapi_client.call(
            "ak.wwise.core.object.get",
            args={"from": {"id": [object_id]}},
            options={"return": ["id", "name", "@OutputBus", "OutputBus", "@VolumeOffset"]},
        )
        volume_offset = 0.0
        if isinstance(obj_result, dict):
            obj_items = obj_result.get("return", [])
            if obj_items:
                effective_bus_id = self._extract_reference_id(
                    obj_items[0].get("OutputBus")
                ) or self._extract_reference_id(obj_items[0].get("@OutputBus"))
                volume_offset = self._safe_db_value(obj_items[0].get("@VolumeOffset"))

        # --- query ancestors ------------------------------------
        result = self.waapi_client.call(
            "ak.wwise.core.object.get",
            args={
                "from": {"id": [object_id]},
                "transform": [{"select": ["ancestors"]}],
            },
            options={"return": ["id", "name", "type", "@Volume", "@MakeUpGain", "@OutputBus"]},
        )
        error_text = self._extract_waapi_error(result)
        if error_text:
            warnings.append(f"{object_id}: 查询祖先层级失败: {error_text}")
        else:
            items = result.get("return", []) if isinstance(result, dict) else []
            for ancestor in (items if isinstance(items, list) else []):
                vol = self._safe_db_value(ancestor.get("@Volume"))
                mug = self._safe_db_value(ancestor.get("@MakeUpGain"))
                total_volume_db += vol
                total_makeup_db += mug
                ancestors_list.append({
                    "id": ancestor.get("id", ""),
                    "name": ancestor.get("name", ""),
                    "type": ancestor.get("type", ""),
                    "volume_db": round(vol, 3),
                    "makeup_gain_db": round(mug, 3),
                })
                # Fallback for cases where the object did not resolve an effective OutputBus.
                if not effective_bus_id:
                    anc_bus = self._extract_reference_id(ancestor.get("@OutputBus"))
                    if anc_bus:
                        effective_bus_id = anc_bus

        return {
            "volume_offset_db": round(volume_offset, 3),
            "ancestor_volume_db": round(total_volume_db, 3),
            "ancestor_makeup_db": round(total_makeup_db, 3),
            "ancestor_gain_db": round(total_volume_db + total_makeup_db, 3),
            "effective_bus_id": effective_bus_id,
            "ancestors": ancestors_list,
            "warnings": warnings,
        }

    def _estimate_bus_gain_chain(self, bus_id: str | None, max_depth: int = 12):
        """Follow the Bus chain starting from *bus_id* and sum @BusVolume + @Volume."""
        bus_nodes: list[dict] = []
        total_bus_volume_db = 0.0
        total_volume_db = 0.0
        visited: set[str] = set()
        current_id = bus_id
        warnings: list[str] = []

        while current_id and len(bus_nodes) < max_depth:
            if current_id in visited:
                bus_nodes.append({"id": current_id, "name": "[CycleDetected]", "type": "Unknown",
                                  "bus_volume_db": 0.0, "volume_db": 0.0})
                warnings.append(f"{current_id}: 检测到 Bus 路由环路，已停止继续追踪。")
                break
            visited.add(current_id)

            fetch_result = self._fetch_object_route_node(current_id)
            warning_text = fetch_result.get("warning")
            if warning_text:
                warnings.append(warning_text)
            node = fetch_result.get("node")
            if not node:
                break

            node_bus_volume = self._safe_db_value(node.get("@BusVolume"))
            node_volume = self._safe_db_value(node.get("@Volume"))
            total_bus_volume_db += node_bus_volume
            total_volume_db += node_volume
            bus_nodes.append({
                "id": node.get("id", current_id),
                "name": node.get("name", ""),
                "type": node.get("type", ""),
                "path": node.get("path", ""),
                "bus_volume_db": round(node_bus_volume, 3),
                "volume_db": round(node_volume, 3),
            })

            # Follow the effective OutputBus first. If a Bus/AuxBus node has no local
            # OutputBus override, route continues through its parent in the hierarchy.
            next_id = self._extract_reference_id(node.get("OutputBus")) or self._extract_reference_id(node.get("@OutputBus"))
            if not next_id:
                node_type = (node.get("type") or "").strip().casefold()
                if node_type in {"bus", "auxbus"}:
                    next_id = self._extract_reference_id(node.get("parent"))
            if not next_id:
                break
            current_id = next_id

        if current_id and len(bus_nodes) >= max_depth:
            warnings.append(f"Bus 路由深度超过 {max_depth} 层，已停止继续追踪。")

        return {
            "bus_volume_db": round(total_bus_volume_db, 3),
            "bus_chain_volume_db": round(total_volume_db, 3),
            "bus_gain_db": round(total_bus_volume_db + total_volume_db, 3),
            "bus_nodes": bus_nodes,
            "warnings": warnings,
        }

    def _estimate_route_gain_chain(self, start_object_id: str, max_depth: int = 12):
        """Full route gain: VolumeOffset + ancestor hierarchy + bus chain."""
        warnings: list[str] = []

        ancestor_summary = self._estimate_ancestor_gain_chain(start_object_id)
        warnings.extend(ancestor_summary.get("warnings", []))

        bus_summary = self._estimate_bus_gain_chain(
            ancestor_summary.get("effective_bus_id"), max_depth=max_depth,
        )
        warnings.extend(bus_summary.get("warnings", []))

        volume_offset = ancestor_summary.get("volume_offset_db", 0.0)
        ancestor_gain = ancestor_summary.get("ancestor_gain_db", 0.0)
        bus_gain = bus_summary.get("bus_gain_db", 0.0)
        route_gain_db = round(volume_offset + ancestor_gain + bus_gain, 3)

        # Build combined route_nodes for backward compatibility
        route_nodes = []
        for anc in ancestor_summary.get("ancestors", []):
            route_nodes.append({
                "id": anc.get("id", ""),
                "name": anc.get("name", ""),
                "type": anc.get("type", ""),
                "volume_db": anc.get("volume_db", 0.0),
                "makeup_gain_db": anc.get("makeup_gain_db", 0.0),
                "section": "actor-mixer",
            })
        for bn in bus_summary.get("bus_nodes", []):
            route_nodes.append({
                "id": bn.get("id", ""),
                "name": bn.get("name", ""),
                "type": bn.get("type", ""),
                "path": bn.get("path", ""),
                "bus_volume_db": bn.get("bus_volume_db", 0.0),
                "volume_db": bn.get("volume_db", 0.0),
                "section": "bus",
            })

        return {
            "route_gain_db": route_gain_db,
            "volume_offset_db": volume_offset,
            "ancestor_gain_db": ancestor_gain,
            "bus_gain_db": bus_gain,
            "route_nodes": route_nodes,
            "ancestor_summary": ancestor_summary,
            "bus_summary": bus_summary,
            "warnings": warnings,
        }

    def analyze_selected_sources_full_route_loudness(self, source_files=None):
        """Estimate full-route loudness for selected objects.

        Method (aligned with get_audio_sources approach):
        1) Resolve AudioFileSource entries from the selection
           (via get_selected_source_files, which traverses descendants,
           filters AudioFileSource, and checks ActiveSource).
        2) For each resolved source entry, analyze source loudness (LUFS-I).
        3) Estimate the full route gain chain starting from the Sound object:
           VolumeOffset + ancestor @Volume/@MakeUpGain + bus @BusVolume/@Volume.
        4) estimated_full_route_lufs = source LUFS-I + total_route_gain.
        """
        warnings: list[str] = []

        # 1. Resolve source files at AudioFileSource level
        if source_files is None:
            source_entries = self.get_selected_source_files()
        else:
            source_entries = source_files

        if not source_entries:
            return AnalysisReport({
                "results": [],
                "warnings": ["No audio source files found in the current selection."],
                "count": 0,
            })

        results = []
        total = len(source_entries)
        try:
            for index, entry in enumerate(source_entries, start=1):
                if not isinstance(entry, dict):
                    continue

                source_path = (entry.get("originalFilePath") or entry.get("path") or "").strip()
                object_name = (
                    entry.get("objectName") or entry.get("name") or source_path or "unknown"
                ).strip()
                object_id = (entry.get("objectId") or "").strip()

                self._report_analysis_progress(index, total, object_name)

                if not source_path or not os.path.isfile(source_path):
                    warnings.append(f"{object_name}: 原始音频文件在本机不存在: {source_path}")
                    continue

                # 2. Analyze source loudness
                try:
                    source_analysis = self.analyze_audio_file(source_path)
                except Exception as exc:
                    warnings.append(f"{object_name}: analyze_audio_file failed: {exc}")
                    continue

                source_lufs = source_analysis.get("integrated_loudness_lufs")

                # 3. Estimate full route gain chain from the Sound object
                if object_id:
                    route_summary = self._estimate_route_gain_chain(object_id)
                else:
                    route_summary = {
                        "route_gain_db": 0.0, "volume_offset_db": 0.0,
                        "ancestor_gain_db": 0.0, "bus_gain_db": 0.0,
                        "route_nodes": [], "ancestor_summary": {}, "bus_summary": {},
                        "warnings": [f"无 objectId，无法估算路由增益。"],
                    }

                route_gain_db = route_summary.get("route_gain_db", 0.0)
                volume_offset_db = route_summary.get("volume_offset_db", 0.0)
                ancestor_gain_db = route_summary.get("ancestor_gain_db", 0.0)
                bus_gain_db = route_summary.get("bus_gain_db", 0.0)
                route_warnings = route_summary.get("warnings", [])
                for warning_text in route_warnings:
                    warnings.append(f"{object_name}: {warning_text}")

                # 4. estimated_full_route_lufs = source LUFS-I + total_route_gain
                estimated_full_route_lufs = None
                if source_lufs is not None:
                    estimated_full_route_lufs = round(float(source_lufs) + float(route_gain_db), 2)

                ancestor_summary = route_summary.get("ancestor_summary", {})
                bus_summary = route_summary.get("bus_summary", {})

                results.append(
                    {
                        "object": {
                            "id": object_id,
                            "name": object_name,
                            "type": entry.get("objectType", ""),
                            "path": entry.get("objectPath", ""),
                        },
                        "source": {
                            "path": source_path,
                            "lufs_i": source_lufs,
                            "analysis": source_analysis,
                        },
                        "route": {
                            "gain_db": route_gain_db,
                            "volume_offset_db": volume_offset_db,
                            "ancestor_gain_db": ancestor_gain_db,
                            "bus_gain_db": bus_gain_db,
                            "nodes": route_summary.get("route_nodes", []),
                            "ancestor_hierarchy": ancestor_summary.get("ancestors", []),
                            "bus_chain": bus_summary.get("bus_nodes", []),
                            "warnings": route_warnings,
                        },
                        "estimated_full_route_lufs": estimated_full_route_lufs,
                    }
                )
        finally:
            self._finish_analysis_progress()

        return AnalysisReport({
            "results": results,
            "warnings": warnings,
            "count": len(results),
        })

    def analyze_project_source_files_loudness(self, limit: int | None = None, source_files=None):
        if source_files is None:
            source_files = self.get_project_source_files()
        if isinstance(limit, int) and limit > 0:
            source_files = source_files[:limit]

        results = []
        warnings = []
        total = len(source_files)
        try:
            for index, item in enumerate(source_files, start=1):
                if not isinstance(item, dict):
                    continue
                source_path = (item.get("originalFilePath") or item.get("path") or "").strip()
                object_name = (item.get("objectName") or item.get("name") or source_path or "unknown").strip()
                self._report_analysis_progress(index, total, object_name)
                if not source_path or not os.path.isfile(source_path):
                    warnings.append(f"{object_name}: 本地文件不存在，已跳过。")
                    continue
                try:
                    analysis = self.analyze_audio_file(source_path)
                except Exception as exc:
                    warnings.append(f"{object_name}: 响度分析失败: {exc}")
                    continue
                results.append(
                    {
                        "objectName": object_name,
                        "objectId": item.get("objectId", ""),
                        "objectPath": item.get("objectPath", ""),
                        "originalFilePath": source_path,
                        "analysis": analysis,
                    }
                )
        finally:
            self._finish_analysis_progress()

        return AnalysisReport({
            "count": len(results),
            "results": results,
            "warnings": warnings,
        })

    @staticmethod
    def _choose_backup_path(normalized: str) -> str:
        stem, ext = os.path.splitext(normalized)
        backup_path = f"{stem}.bak{ext}"
        counter = 1
        while os.path.exists(backup_path):
            backup_path = f"{stem}.bak{counter}{ext}"
            counter += 1
        return backup_path

    def _atomic_write_audio(self, normalized, audio_data, sample_rate, *, subtype, fmt,
                            backup, sf):
        """Crash-safe audio write: write a temp file in the same dir, fsync,
        os.replace() onto the target, and only then move the original to *.bak.

        Order matters: the new data is fully written and atomically swapped in
        BEFORE the original is touched, so an interrupt at any point leaves
        either the old file or the new file at ``normalized`` — never a hole.
        Returns the backup path (or None).
        """
        directory = os.path.dirname(normalized) or "."
        backup_path = self._choose_backup_path(normalized) if backup else None

        # 1) Back up the original first (copy, not move) so it survives a crash.
        if backup:
            shutil.copy2(normalized, backup_path)

        # 2) Write the new audio to a temp file in the same directory, fsync it.
        fd, tmp_path = tempfile.mkstemp(suffix=os.path.splitext(normalized)[1] or ".tmp",
                                        dir=directory)
        os.close(fd)
        try:
            sf.write(tmp_path, audio_data, sample_rate, subtype=subtype, format=fmt)
            # Best-effort fsync so the bytes are durable before the swap.
            try:
                with open(tmp_path, "rb") as fh:
                    os.fsync(fh.fileno())
            except OSError:
                pass
            # 3) Atomic swap onto the target.
            os.replace(tmp_path, normalized)
        except Exception:
            # Write failed: clean up temp, and since the original is untouched
            # (we only copied it), the target still holds the original bytes.
            if os.path.exists(tmp_path):
                try:
                    os.remove(tmp_path)
                except OSError:
                    pass
            # The redundant backup copy is now pointless; remove it to avoid clutter.
            if backup and backup_path and os.path.exists(backup_path):
                try:
                    os.remove(backup_path)
                except OSError:
                    pass
            raise
        return backup_path

    def normalize_audio_loudness(
        self,
        path: str,
        target_lufs: float = -16.0,
        backup: bool = True,
        *,
        stage: bool = True,
    ) -> dict:
        """Normalize the integrated loudness of an audio file to target_lufs.

        Args:
            path: Path to the audio file (WAV, FLAC, OGG, etc.).
            target_lufs: Target integrated loudness in LUFS (default -16).
            backup: If True (default), keep the original as a *.bak.* copy.
            stage: If True (default) and a ``file_write_stager`` is wired, the
                disk write is deferred to the GUI confirmation pipeline instead
                of applying immediately. Pass stage=False to force an immediate
                (still atomic) write — used by batch apply, which confirms once
                up front.

        The write is crash-safe: the new audio is written to a temp file and
        atomically swapped in (os.replace) only after the original is backed up,
        so a failure never leaves the target path empty.

        Returns:
            dict with path, original_lufs, target_lufs, result_lufs, backup_path.
            When the write is staged, ``pending_confirmation`` is True and
            ``result_lufs`` is the predicted (pre-write) value.
        """
        normalized = self._normalize_path(path)
        if not os.path.exists(normalized):
            raise FileNotFoundError(normalized)
        if os.path.isdir(normalized):
            raise IsADirectoryError(normalized)

        self.remember_paths([normalized], source="local")

        np, _librosa, pyln, sf = self._import_audio_analysis_dependencies()

        # Read with soundfile (preserves format metadata)
        info = sf.info(normalized)
        audio, sample_rate = sf.read(normalized, dtype="float64", always_2d=True)
        if audio.size == 0:
            raise ValueError("Audio file is empty or could not be decoded.")

        meter_input = self._coerce_meter_input(audio, np)

        # Measure original loudness (with short-audio fallback)
        original_lufs, _, norm_warnings = self._measure_loudness_with_fallback(
            meter_input, sample_rate, pyln, np,
        )

        if not self._is_finite_number(original_lufs):
            raise ValueError(
                f"Cannot measure loudness (got {original_lufs}). "
                "The file may be silent or too short."
            )

        original_lufs = float(original_lufs)

        # Normalize (in memory)
        loudness_normalized = pyln.normalize.loudness(
            audio, original_lufs, target_lufs,
        )

        # Clip to prevent clipping artifacts on int formats
        peak = float(np.max(np.abs(loudness_normalized)))
        if peak > 1.0:
            loudness_normalized = loudness_normalized / peak

        write_subtype = getattr(info, "subtype", None)
        write_format = getattr(info, "format", None)

        def _apply_and_verify() -> dict:
            backup_path = self._atomic_write_audio(
                normalized, loudness_normalized, sample_rate,
                subtype=write_subtype, fmt=write_format, backup=backup, sf=sf,
            )
            # Verify result (with short-audio fallback)
            verify_audio, _ = sf.read(normalized, dtype="float64", always_2d=True)
            verify_input = self._coerce_meter_input(verify_audio, np)
            result_lufs, _, _ = self._measure_loudness_with_fallback(
                verify_input, sample_rate, pyln, np,
            )
            return {
                "path": normalized,
                "original_lufs": round(original_lufs, 2),
                "target_lufs": round(target_lufs, 2),
                "result_lufs": round(float(result_lufs), 2) if self._is_finite_number(result_lufs) else None,
                "backup_path": backup_path,
                "peak_after_normalize": round(peak, 4),
                "normalization_warnings": norm_warnings,
            }

        # Defer to the GUI confirmation pipeline when available and requested.
        if stage and callable(self.file_write_stager):
            self.file_write_stager(normalized, _apply_and_verify)
            return {
                "path": normalized,
                "original_lufs": round(original_lufs, 2),
                "target_lufs": round(target_lufs, 2),
                "result_lufs": round(float(target_lufs), 2),  # predicted
                "backup_path": None,
                "peak_after_normalize": round(peak, 4),
                "normalization_warnings": norm_warnings,
                "pending_confirmation": True,
            }

        return _apply_and_verify()

    # -----------------------------------------------------------------
    # Phase 1: loudness compliance health-check + batch fix
    # -----------------------------------------------------------------

    @staticmethod
    def _is_backup_audio_file(path: str) -> bool:
        """True for files normalize_audio_loudness created as backups
        (``name.bak.wav`` / ``name.bak1.wav``). They must be excluded from
        batch scans so re-runs don't normalize backups into .bak.bak chains."""
        stem = os.path.splitext(os.path.basename(str(path or "")))[0]
        tail = stem.rsplit(".", 1)[-1].lower() if "." in stem else ""
        return tail == "bak" or (tail.startswith("bak") and tail[3:].isdigit())

    def _evaluate_loudness_compliance(self, row: dict, lo: float, hi: float, true_peak_limit_dbfs: float) -> dict:
        """Annotate one analysis row with a pass/fail verdict against a target
        LUFS range and a true-peak limit. Returns the same row (mutated)."""
        lufs = row.get("integrated_loudness_lufs")
        issues: list[str] = []
        loudness_deviation = 0.0
        true_peak_over = None

        if not self._is_finite_number(lufs):
            row["compliance"] = "unknown"
            row["compliance_issues"] = ["loudness_unknown"]
            row["loudness_deviation_db"] = None
            row["true_peak_over_db"] = None
            return row

        lufs = float(lufs)
        if lufs < lo:
            issues.append("loudness_low")
            loudness_deviation = lo - lufs
        elif lufs > hi:
            issues.append("loudness_high")
            loudness_deviation = lufs - hi

        tp = row.get("true_peak_dbfs")
        if self._is_finite_number(tp) and float(tp) > true_peak_limit_dbfs:
            issues.append("true_peak_over")
            true_peak_over = float(tp) - true_peak_limit_dbfs

        row["compliance"] = "pass" if not issues else "fail"
        row["compliance_issues"] = issues
        row["loudness_deviation_db"] = self._round_or_none(loudness_deviation)
        row["true_peak_over_db"] = self._round_or_none(true_peak_over)
        return row

    def check_directory_loudness_compliance(
        self,
        path: str,
        *,
        target_lufs_min: float = -16.0,
        target_lufs_max: float = -12.0,
        true_peak_limit_dbfs: float = -1.0,
        recursive: bool = True,
        extensions: list[str] | None = None,
        limit: int | None = None,
    ):
        """Analyse a folder of audio and flag files outside a loudness range /
        over a true-peak limit. Reuses ``analyze_directory_loudness`` and adds a
        per-file pass/fail verdict plus a compliance summary. Read-only.

        Returns an ``AnalysisReport`` ({count, results, warnings, summary}) where
        ``summary['compliance']`` holds counts and ``noncompliant_files`` is the
        worst-first list the GUI/LLM can act on.
        """
        report = self.analyze_directory_loudness(
            path,
            recursive=recursive,
            extensions=extensions,
            limit=limit,
        )
        lo = min(float(target_lufs_min), float(target_lufs_max))
        hi = max(float(target_lufs_min), float(target_lufs_max))

        # Exclude backup siblings (*.bak.wav) so a health-check / fix never
        # re-processes files a previous normalize created.
        all_results = report.get("results", [])
        results = [r for r in all_results if not self._is_backup_audio_file(r.get("path") or r.get("file") or "")]
        report["results"] = results
        report["count"] = len(results)
        compliant = 0
        unknown = 0
        noncompliant: list[dict] = []
        for row in results:
            self._evaluate_loudness_compliance(row, lo, hi, float(true_peak_limit_dbfs))
            verdict = row.get("compliance")
            if verdict == "pass":
                compliant += 1
            elif verdict == "unknown":
                unknown += 1
            else:
                noncompliant.append(row)

        def _worst(row: dict) -> float:
            return max(
                float(row.get("loudness_deviation_db") or 0.0),
                float(row.get("true_peak_over_db") or 0.0),
            )

        noncompliant.sort(key=_worst, reverse=True)

        summary = dict(report.get("summary", {}))
        summary["compliance"] = {
            "target_lufs_min": round(lo, 2),
            "target_lufs_max": round(hi, 2),
            "true_peak_limit_dbfs": round(float(true_peak_limit_dbfs), 2),
            "compliant_count": compliant,
            "noncompliant_count": len(noncompliant),
            "unknown_count": unknown,
            "noncompliant_files": [
                {
                    "file": r.get("file") or r.get("file_name"),
                    "path": r.get("path"),
                    "integrated_loudness_lufs": r.get("integrated_loudness_lufs"),
                    "true_peak_dbfs": r.get("true_peak_dbfs"),
                    "issues": r.get("compliance_issues", []),
                    "loudness_deviation_db": r.get("loudness_deviation_db"),
                    "true_peak_over_db": r.get("true_peak_over_db"),
                }
                for r in noncompliant
            ],
        }
        report["summary"] = summary
        return report

    def batch_normalize_directory_to_target(
        self,
        path: str,
        target_lufs: float = -16.0,
        *,
        target_lufs_min: float | None = None,
        target_lufs_max: float | None = None,
        true_peak_limit_dbfs: float = -1.0,
        only_noncompliant: bool = True,
        recursive: bool = True,
        extensions: list[str] | None = None,
        limit: int | None = None,
        backup: bool = True,
        apply: bool = False,
    ):
        """Batch-normalize a folder to a target LUFS, report-then-apply.

        By default (``apply=False``) this is a DRY RUN: it returns the plan — which
        files would be changed, their current vs target loudness — and writes
        nothing. Call again with ``apply=True`` to perform the (irreversible)
        normalization. This keeps the destructive disk writes behind an explicit
        second confirmation rather than firing on the first call.

        ``only_noncompliant`` (default) skips files already inside
        [target_lufs_min, target_lufs_max] (defaults to a +/-0.5 LU window around
        ``target_lufs`` when the explicit bounds are not given).
        """
        if target_lufs_min is None or target_lufs_max is None:
            lo = float(target_lufs) - 0.5
            hi = float(target_lufs) + 0.5
        else:
            lo = min(float(target_lufs_min), float(target_lufs_max))
            hi = max(float(target_lufs_min), float(target_lufs_max))

        compliance = self.check_directory_loudness_compliance(
            path,
            target_lufs_min=lo,
            target_lufs_max=hi,
            true_peak_limit_dbfs=true_peak_limit_dbfs,
            recursive=recursive,
            extensions=extensions,
            limit=limit,
        )
        results = compliance.get("results", [])

        planned: list[dict] = []
        skipped: list[dict] = []
        for row in results:
            file_path = row.get("path")
            verdict = row.get("compliance")
            current_lufs = row.get("integrated_loudness_lufs")
            if verdict == "unknown":
                skipped.append({"path": file_path, "reason": "无法测量响度，已跳过。"})
                continue
            if only_noncompliant and verdict == "pass":
                skipped.append({"path": file_path, "reason": "已符合目标范围，已跳过。"})
                continue
            gain_db = None
            if self._is_finite_number(current_lufs):
                gain_db = self._round_or_none(float(target_lufs) - float(current_lufs))
            planned.append({
                "path": file_path,
                "current_lufs": current_lufs,
                "target_lufs": round(float(target_lufs), 2),
                "gain_db": gain_db,
                "issues": row.get("compliance_issues", []),
            })

        if not apply:
            return AnalysisReport({
                "mode": "dry_run",
                "count": len(planned),
                "results": planned,
                "skipped": skipped,
                "warnings": [
                    "这是预演（dry run），未修改任何文件。确认无误后，请以 apply=True 再次调用以实际归一化。",
                ],
                "summary": {
                    "planned_count": len(planned),
                    "skipped_count": len(skipped),
                    "target_lufs": round(float(target_lufs), 2),
                },
            })

        applied: list[dict] = []
        failed: list[dict] = []
        total = len(planned)
        try:
            for index, item in enumerate(planned, start=1):
                file_path = item["path"]
                self._report_analysis_progress(index, total, os.path.basename(file_path or ""))
                try:
                    res = self.normalize_audio_loudness(
                        file_path,
                        target_lufs=float(target_lufs),
                        backup=backup,
                        stage=False,  # batch apply confirms once up front; write now (atomically)
                    )
                    applied.append(res)
                except Exception as exc:
                    failed.append({"path": file_path, "error": str(exc)})
        finally:
            self._finish_analysis_progress()

        return AnalysisReport({
            "mode": "applied",
            "count": len(applied),
            "results": applied,
            "failed": failed,
            "skipped": skipped,
            "warnings": [f"{item['path']}: 归一化失败: {item['error']}" for item in failed],
            "summary": {
                "applied_count": len(applied),
                "failed_count": len(failed),
                "skipped_count": len(skipped),
                "target_lufs": round(float(target_lufs), 2),
            },
        })

    # -----------------------------------------------------------------
    # Phase 2: silence / clipping / anomaly detection
    # -----------------------------------------------------------------

    # Default thresholds (overridable per call).
    CLIP_THRESHOLD = 0.999            # |sample| at/above this counts as clipped
    CLIP_MIN_RUN = 3                  # consecutive clipped samples to flag a clip run
    SILENCE_RMS_DBFS = -60.0          # whole-file RMS below this -> (near) silent
    DC_OFFSET_WARN = 0.01             # |mean| above this (linear) -> DC offset
    TRUE_PEAK_OVER_DBFS = 0.0         # inter-sample peak above this -> over (dBTP)
    SHORT_DURATION_SECONDS = 0.1      # shorter than this -> too-short flag

    def _detect_clipping(self, audio, np, *, threshold: float, min_run: int) -> dict:
        """Per-channel consecutive-sample clip-run detection on the raw array.

        Returns {clipped_sample_count, clip_run_count, max_clip_run}. A single
        global max(abs) cannot tell a momentary 0 dBFS sample from sustained
        clipping, so we count runs of |sample| >= threshold of length >= min_run.
        """
        data = np.asarray(audio, dtype=np.float64)
        if data.ndim == 1:
            data = data.reshape(-1, 1)
        clipped_total = 0
        run_count = 0
        max_run = 0
        for ch in range(data.shape[1]):
            mask = np.abs(data[:, ch]) >= threshold
            clipped_total += int(np.count_nonzero(mask))
            if not mask.any():
                continue
            # Run-length encode the boolean mask.
            idx = np.flatnonzero(np.diff(mask.astype(np.int8)))
            bounds = np.concatenate(([0], idx + 1, [mask.size]))
            for start, end in zip(bounds[:-1], bounds[1:]):
                if mask[start]:
                    run_len = int(end - start)
                    if run_len >= min_run:
                        run_count += 1
                        if run_len > max_run:
                            max_run = run_len
        return {
            "clipped_sample_count": clipped_total,
            "clip_run_count": run_count,
            "max_clip_run": max_run,
        }

    def _detect_dc_offset(self, audio, np) -> float:
        """Max absolute per-channel mean (DC offset) on the raw array."""
        data = np.asarray(audio, dtype=np.float64)
        if data.ndim == 1:
            data = data.reshape(-1, 1)
        if data.shape[0] == 0:
            return 0.0
        return float(np.max(np.abs(np.mean(data, axis=0))))

    def detect_audio_anomalies(
        self,
        path: str,
        *,
        clip_threshold: float | None = None,
        clip_min_run: int | None = None,
        silence_rms_dbfs: float | None = None,
        dc_offset_warn: float | None = None,
        true_peak_over_dbfs: float | None = None,
        expected_sample_rates: list[int] | None = None,
        expected_channels: list[int] | None = None,
    ) -> dict:
        """Detect common defects in one audio file: clipping, DC offset, silence,
        true-peak overs, too-short, and (optionally) abnormal sample rate / channel
        count. Runs detection BEFORE any spectral/LUFS analysis so empty, silent or
        ultra-short files still get a verdict instead of raising.

        Returns a dict with raw metrics plus ``anomalies`` (list of codes) and
        ``has_anomaly`` (bool).
        """
        clip_threshold = self.CLIP_THRESHOLD if clip_threshold is None else float(clip_threshold)
        clip_min_run = self.CLIP_MIN_RUN if clip_min_run is None else int(clip_min_run)
        silence_rms_dbfs = self.SILENCE_RMS_DBFS if silence_rms_dbfs is None else float(silence_rms_dbfs)
        dc_offset_warn = self.DC_OFFSET_WARN if dc_offset_warn is None else float(dc_offset_warn)
        true_peak_over_dbfs = self.TRUE_PEAK_OVER_DBFS if true_peak_over_dbfs is None else float(true_peak_over_dbfs)

        normalized = self._normalize_path(path)
        if not os.path.exists(normalized):
            raise FileNotFoundError(normalized)
        if os.path.isdir(normalized):
            raise IsADirectoryError(normalized)

        self.remember_paths([normalized], source="local")
        np, librosa, _pyln, sf = self._import_audio_analysis_dependencies()
        audio, sample_rate, source_info = self._load_audio_samples(normalized, np, librosa, sf)

        frame_count = int(audio.shape[0]) if audio.ndim else 0
        channels = int(audio.shape[1]) if audio.ndim > 1 else 1
        duration_seconds = frame_count / float(sample_rate) if sample_rate else 0.0

        anomalies: list[str] = []
        result: dict = {
            "path": normalized,
            "format": source_info.get("format", "unknown"),
            "subtype": source_info.get("subtype", "unknown"),
            "channels": channels,
            "sample_rate": int(sample_rate) if sample_rate else 0,
            "frame_count": frame_count,
            "duration_seconds": round(duration_seconds, 4),
        }

        # Empty / undecodable.
        if audio.size == 0 or frame_count == 0:
            anomalies.append("empty")
            result["anomalies"] = anomalies
            result["has_anomaly"] = True
            return result

        peak = float(np.max(np.abs(audio)))
        rms = float(np.sqrt(np.mean(np.square(audio))))
        peak_dbfs = self._linear_to_dbfs(peak)
        rms_dbfs = self._linear_to_dbfs(rms)
        true_peak_dbfs = self._compute_true_peak_dbfs(audio, sample_rate, np, librosa)
        if true_peak_dbfs is None:
            true_peak_dbfs = peak_dbfs
        dc_offset = self._detect_dc_offset(audio, np)
        clip = self._detect_clipping(audio, np, threshold=clip_threshold, min_run=clip_min_run)

        result.update({
            "peak_dbfs": peak_dbfs,
            "rms_dbfs": rms_dbfs,
            "true_peak_dbfs": true_peak_dbfs,
            "dc_offset": round(dc_offset, 6),
            **clip,
        })

        # --- anomaly rules ---
        if rms <= 0 or (rms_dbfs is not None and rms_dbfs < silence_rms_dbfs):
            anomalies.append("silent")
        if clip["clip_run_count"] > 0:
            anomalies.append("clipping")
        if self._is_finite_number(true_peak_dbfs) and float(true_peak_dbfs) > true_peak_over_dbfs:
            anomalies.append("true_peak_over")
        if dc_offset > dc_offset_warn:
            anomalies.append("dc_offset")
        if 0 < duration_seconds < self.SHORT_DURATION_SECONDS:
            anomalies.append("too_short")
        if expected_sample_rates and sample_rate and int(sample_rate) not in {int(s) for s in expected_sample_rates}:
            anomalies.append("abnormal_sample_rate")
        if expected_channels and channels not in {int(c) for c in expected_channels}:
            anomalies.append("abnormal_channel_count")

        result["anomalies"] = anomalies
        result["has_anomaly"] = bool(anomalies)
        return result

    def detect_directory_anomalies(
        self,
        path: str,
        *,
        recursive: bool = True,
        extensions: list[str] | None = None,
        limit: int | None = None,
        clip_threshold: float | None = None,
        clip_min_run: int | None = None,
        silence_rms_dbfs: float | None = None,
        dc_offset_warn: float | None = None,
        true_peak_over_dbfs: float | None = None,
        expected_sample_rates: list[int] | None = None,
        expected_channels: list[int] | None = None,
        use_rules: bool = True,
    ):
        """Batch anomaly scan over a folder. Lightweight relative to a full
        loudness pass — skips LUFS/spectrum. Returns an ``AnalysisReport``
        ({count, results, warnings, summary}); ``results`` holds only the files
        that have at least one anomaly, with a per-code tally in ``summary``.

        When ``use_rules`` is True and the caller does not pass explicit
        ``expected_sample_rates`` / ``expected_channels``, the values are taken
        from the audio rules config (Phase 3), so a single config drives both
        structure validation and the abnormal-sr/channel anomaly checks."""
        if use_rules and (expected_sample_rates is None or expected_channels is None):
            audio_rules = (self._load_audio_rules().get("audio") or {})
            if expected_sample_rates is None:
                expected_sample_rates = audio_rules.get("expected_sample_rates") or None
            if expected_channels is None:
                expected_channels = audio_rules.get("expected_channels") or None
        root = self._normalize_path(path)
        if not os.path.exists(root):
            raise FileNotFoundError(root)
        normalized_exts = {
            str(ext).lower() if str(ext).startswith(".") else f".{str(ext).lower()}"
            for ext in (extensions or [".wav"])
        }
        if os.path.isfile(root):
            files = [root]
            root_path = os.path.dirname(root)
        elif os.path.isdir(root):
            root_path = root
            files = []
            if recursive:
                for current_root, _dirs, names in os.walk(root):
                    for name in names:
                        if os.path.splitext(name)[1].lower() in normalized_exts:
                            files.append(os.path.join(current_root, name))
            else:
                for name in os.listdir(root):
                    full = os.path.join(root, name)
                    if os.path.isfile(full) and os.path.splitext(name)[1].lower() in normalized_exts:
                        files.append(full)
            files = sorted(files, key=str.lower)
        else:
            raise ValueError(f"Unsupported local path: {root}")

        files = [f for f in files if not self._is_backup_audio_file(f)]
        if isinstance(limit, int) and limit > 0:
            files = files[:limit]

        flagged: list[dict] = []
        warnings: list[str] = []
        tally: dict[str, int] = {}
        total = len(files)
        scanned = 0
        try:
            for index, file_path in enumerate(files, start=1):
                self._report_analysis_progress(index, total, os.path.basename(file_path))
                try:
                    detection = self.detect_audio_anomalies(
                        file_path,
                        clip_threshold=clip_threshold,
                        clip_min_run=clip_min_run,
                        silence_rms_dbfs=silence_rms_dbfs,
                        dc_offset_warn=dc_offset_warn,
                        true_peak_over_dbfs=true_peak_over_dbfs,
                        expected_sample_rates=expected_sample_rates,
                        expected_channels=expected_channels,
                    )
                except Exception as exc:
                    warnings.append(f"{file_path}: 异常检测失败: {exc}")
                    continue
                scanned += 1
                if detection.get("has_anomaly"):
                    detection["file"] = os.path.basename(file_path)
                    detection["file_name"] = os.path.basename(file_path)
                    detection["relative_path"] = os.path.relpath(file_path, root_path)
                    flagged.append(detection)
                    for code in detection.get("anomalies", []):
                        tally[code] = tally.get(code, 0) + 1
        finally:
            self._finish_analysis_progress()

        if total == 0:
            warnings.append("未找到匹配的音频文件。")

        self.remember_paths([root_path] + files, source="local")
        return AnalysisReport({
            "root_path": root_path,
            "file_count": total,
            "scanned_count": scanned,
            "count": len(flagged),
            "results": flagged,
            "warnings": warnings,
            "summary": {
                "scanned_count": scanned,
                "flagged_count": len(flagged),
                "anomaly_tally": tally,
            },
        })

    # -----------------------------------------------------------------
    # Phase 3: naming / project-structure validation (config-driven)
    # -----------------------------------------------------------------

    # Container types whose emptiness is worth flagging. Excludes Sound/leaf
    # types (a Sound legitimately has no child objects).
    _CONTAINER_TYPES = {
        "randomsequencecontainer", "switchcontainer", "blendcontainer",
        "actormixer", "folder", "workunit", "virtualfolder",
    }

    # Built-in default rules. Written to AUDIO_RULES_FILE on first use so the
    # team can edit them. Kept deliberately permissive — naming patterns are
    # empty by default (every name passes) so the feature never spams false
    # positives until a team opts in to a convention.
    DEFAULT_AUDIO_RULES = {
        "naming": {
            # type (lowercased) -> regex the object name must fully match.
            # Empty map = no naming rules enforced.
            "patterns": {},
            # Object types to skip entirely during naming validation.
            "ignore_types": ["workunit", "folder", "virtualfolder"],
        },
        "audio": {
            # Empty list = accept any. Used by detect_* and structure validation.
            "expected_sample_rates": [],
            "expected_channels": [],
        },
        "structure": {
            "flag_empty_containers": True,
            "flag_missing_source_files": True,
        },
    }

    def _audio_rules_path(self):
        try:
            from src.utils.app_paths import AUDIO_RULES_FILE
            return str(AUDIO_RULES_FILE)
        except Exception:
            return os.path.join(os.path.expanduser("~"), ".audiomate_audio_rules.json")

    @staticmethod
    def _deep_merge_defaults(base: dict, override: dict) -> dict:
        """Return base recursively overlaid with override (override wins)."""
        merged = dict(base)
        for key, value in (override or {}).items():
            if isinstance(value, dict) and isinstance(merged.get(key), dict):
                merged[key] = AgentToolbox._deep_merge_defaults(merged[key], value)
            else:
                merged[key] = value
        return merged

    def _load_audio_rules(self, write_default_if_missing: bool = True) -> dict:
        """Load the audio rules config, falling back to (and optionally seeding)
        the built-in defaults. Always returns a complete dict — user files only
        need to specify the keys they want to override."""
        path = self._audio_rules_path()
        if os.path.isfile(path):
            try:
                with open(path, "r", encoding="utf-8") as handle:
                    user_rules = json.load(handle)
                if isinstance(user_rules, dict):
                    return self._deep_merge_defaults(self.DEFAULT_AUDIO_RULES, user_rules)
            except Exception:
                pass  # fall through to defaults on malformed config
        elif write_default_if_missing:
            try:
                os.makedirs(os.path.dirname(path) or ".", exist_ok=True)
                with open(path, "w", encoding="utf-8") as handle:
                    json.dump(self.DEFAULT_AUDIO_RULES, handle, ensure_ascii=False, indent=2)
            except Exception:
                pass
        return dict(self.DEFAULT_AUDIO_RULES)

    def _validate_object_name(self, name: str, obj_type: str, naming_rules: dict) -> str | None:
        """Return a violation reason string if the name breaks a naming rule,
        else None."""
        patterns = naming_rules.get("patterns") or {}
        ignore_types = {str(t).lower() for t in (naming_rules.get("ignore_types") or [])}
        type_key = (obj_type or "").strip().lower()
        if type_key in ignore_types:
            return None
        pattern = patterns.get(type_key)
        if not pattern:
            return None
        try:
            if re.fullmatch(pattern, name or "") is None:
                return f"名称 '{name}' 不符合 {obj_type} 的命名规则 /{pattern}/"
        except re.error:
            return None  # invalid pattern -> don't block
        return None

    def validate_project_structure(
        self,
        *,
        scope: str = "project",
        rules: dict | None = None,
    ):
        """Validate Wwise project structure against config rules: empty containers,
        objects with a set originalFilePath whose local file is missing, and naming
        convention violations. Read-only.

        ``scope`` is ``'project'`` (walk Containers/Actor-Mixer + Interactive Music
        roots) or ``'selection'`` (current selection + descendants). Returns an
        ``AnalysisReport`` with per-issue rows and a tally in ``summary``.
        """
        if self.waapi_client is None or not getattr(self.waapi_client, "connected", False):
            return AnalysisReport({
                "count": 0,
                "results": [],
                "warnings": ["未连接 Wwise，无法校验工程结构。"],
                "summary": {"issue_tally": {}},
            })

        active_rules = rules if isinstance(rules, dict) else self._load_audio_rules()
        naming_rules = active_rules.get("naming", {}) or {}
        structure_rules = active_rules.get("structure", {}) or {}
        flag_empty = structure_rules.get("flag_empty_containers", True)
        flag_missing = structure_rules.get("flag_missing_source_files", True)

        warnings: list[str] = []
        objects = self._collect_objects_for_validation(scope, warnings)

        issues: list[dict] = []
        tally: dict[str, int] = {}

        def _add_issue(obj: dict, code: str, detail: str):
            issues.append({
                "issue": code,
                "object_id": (obj.get("id") or "").strip(),
                "object_name": (obj.get("name") or "").strip(),
                "object_type": (obj.get("type") or "").strip(),
                "object_path": (obj.get("path") or "").strip(),
                "detail": detail,
            })
            tally[code] = tally.get(code, 0) + 1

        for obj in objects:
            if not isinstance(obj, dict):
                continue
            obj_type = (obj.get("type") or "").strip()
            obj_name = (obj.get("name") or "").strip()
            type_key = obj_type.lower()

            # Empty container
            if flag_empty and type_key in self._CONTAINER_TYPES:
                children = obj.get("childrenCount")
                if isinstance(children, (int, float)) and int(children) == 0:
                    _add_issue(obj, "empty_container", "容器为空，没有任何子对象。")

            # Missing source file (originalFilePath set but file absent on disk)
            if flag_missing:
                raw_path = (obj.get("originalFilePath") or "").strip()
                if raw_path:
                    resolved = self._resolve_project_source_path(raw_path)
                    if not resolved or not os.path.isfile(resolved):
                        _add_issue(obj, "missing_source_file", f"源文件在本地缺失：{raw_path}")

            # Naming violation
            violation = self._validate_object_name(obj_name, obj_type, naming_rules)
            if violation:
                _add_issue(obj, "naming_violation", violation)

        return AnalysisReport({
            "scope": scope,
            "count": len(issues),
            "object_count": len(objects),
            "results": issues,
            "warnings": warnings,
            "summary": {
                "object_count": len(objects),
                "issue_count": len(issues),
                "issue_tally": tally,
                "rules_source": "custom" if rules is not None else self._audio_rules_path(),
            },
        })

    _VALIDATION_RETURN_FIELDS = [
        "id", "name", "type", "path", "childrenCount", "originalFilePath",
    ]

    def _collect_objects_for_validation(self, scope: str, warnings: list[str]) -> list[dict]:
        """Gather objects (with childrenCount + originalFilePath) for validation,
        either from the current selection's descendants or the project roots."""
        if str(scope).strip().lower() == "selection":
            all_objects, sel_warnings = self._get_selected_and_descendant_objects()
            warnings.extend(sel_warnings)
            ids = [o.get("id") for o in all_objects if isinstance(o, dict) and o.get("id")]
            if not ids:
                return []
            items, q_warnings = self._query_objects_by_ids_resilient(ids, self._VALIDATION_RETURN_FIELDS)
            warnings.extend(q_warnings)
            return items

        # Project scope: walk descendants of the relevant hierarchy roots.
        roots = self._validation_root_paths()
        collected: list[dict] = []
        seen_ids: set[str] = set()
        for root_path in roots:
            result = self.waapi_client.call(
                "ak.wwise.core.object.get",
                args={"from": {"path": [root_path]}, "transform": [{"select": ["descendants"]}]},
                options={"return": self._VALIDATION_RETURN_FIELDS},
            )
            error_text = self._extract_waapi_error(result)
            if error_text:
                # A root that doesn't exist in this project/version is fine; skip it.
                continue
            items = result.get("return", []) if isinstance(result, dict) else []
            for obj in (items if isinstance(items, list) else []):
                if not isinstance(obj, dict):
                    continue
                obj_id = (obj.get("id") or "").strip()
                if obj_id and obj_id in seen_ids:
                    continue
                if obj_id:
                    seen_ids.add(obj_id)
                collected.append(obj)
        if not collected:
            warnings.append("未能在工程层级根下枚举到对象（可能 Wwise 版本/路径不匹配或工程为空）。")
        return collected

    def _validation_root_paths(self) -> list[str]:
        """Version-aware hierarchy roots to walk for project-scope validation.
        Probes get_version_context when available, else tries both layouts."""
        is_2025 = False
        client = self.waapi_client
        try:
            if hasattr(client, "get_version_context"):
                info = client.get_version_context() or {}
                is_2025 = bool(info.get("is_2025_or_later")) or int(info.get("year") or 0) >= 2025
        except Exception:
            is_2025 = False
        if is_2025:
            return ["\\Containers", "\\Actor-Mixer Hierarchy", "\\Interactive Music Hierarchy"]
        return ["\\Actor-Mixer Hierarchy", "\\Interactive Music Hierarchy", "\\Containers"]


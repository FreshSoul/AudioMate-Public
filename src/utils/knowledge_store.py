"""
知识库存储层 — 在 KNOWLEDGE_DIR 下管理多个知识库，
每个知识库是一个子目录，包含 meta.json 和上传的文档文件。
"""

import json
import os
import re
import shutil
import uuid
from datetime import datetime

from src.utils.app_paths import KNOWLEDGE_DIR as _KNOWLEDGE_DIR

SUPPORTED_KNOWLEDGE_EXTENSIONS = {
    ".txt", ".md", ".csv", ".json", ".xml", ".html", ".log",
    ".pdf", ".docx", ".xlsx", ".pptx",
}

KNOWLEDGE_DIR = str(_KNOWLEDGE_DIR)


def _ensure_dir(path: str):
    os.makedirs(path, exist_ok=True)


def is_supported_knowledge_file(path: str) -> bool:
    return os.path.splitext(path)[1].lower() in SUPPORTED_KNOWLEDGE_EXTENSIONS


def _unique_destination_path(kb_path: str, filename: str) -> str:
    stem, ext = os.path.splitext(filename)
    candidate = os.path.join(kb_path, filename)
    index = 1
    while os.path.exists(candidate):
        candidate = os.path.join(kb_path, f"{stem}_{index}{ext}")
        index += 1
    return candidate


# ── Knowledge‑base CRUD ────────────────────────────────────────────

def list_knowledge_bases() -> list[dict]:
    """返回所有知识库 [{id, name, file_count, total_size, created_at}]"""
    _ensure_dir(KNOWLEDGE_DIR)
    result = []
    for name in os.listdir(KNOWLEDGE_DIR):
        kb_path = os.path.join(KNOWLEDGE_DIR, name)
        meta_path = os.path.join(kb_path, "meta.json")
        if not os.path.isdir(kb_path) or not os.path.exists(meta_path):
            continue
        try:
            with open(meta_path, "r", encoding="utf-8") as f:
                meta = json.load(f)
        except Exception:
            continue
        files = _list_files(kb_path)
        total_size = sum(f["size"] for f in files)
        result.append({
            "id": meta.get("id", name),
            "name": meta.get("name", name),
            "file_count": len(files),
            "total_size": total_size,
            "created_at": meta.get("created_at", ""),
        })
    result.sort(key=lambda x: x["created_at"], reverse=True)
    return result


def create_knowledge_base(name: str) -> str:
    """创建空知识库，返回 kb_id"""
    _ensure_dir(KNOWLEDGE_DIR)
    kb_id = str(uuid.uuid4())[:8]
    kb_path = os.path.join(KNOWLEDGE_DIR, kb_id)
    os.makedirs(kb_path, exist_ok=True)
    meta = {"id": kb_id, "name": name, "created_at": datetime.now().isoformat()}
    with open(os.path.join(kb_path, "meta.json"), "w", encoding="utf-8") as f:
        json.dump(meta, f, ensure_ascii=False, indent=2)
    return kb_id


def rename_knowledge_base(kb_id: str, new_name: str) -> bool:
    kb_path = os.path.join(KNOWLEDGE_DIR, kb_id)
    meta_path = os.path.join(kb_path, "meta.json")
    if not os.path.exists(meta_path):
        return False
    with open(meta_path, "r", encoding="utf-8") as f:
        meta = json.load(f)
    meta["name"] = new_name
    with open(meta_path, "w", encoding="utf-8") as f:
        json.dump(meta, f, ensure_ascii=False, indent=2)
    return True


def delete_knowledge_base(kb_id: str) -> bool:
    kb_path = os.path.join(KNOWLEDGE_DIR, kb_id)
    if os.path.isdir(kb_path):
        shutil.rmtree(kb_path)
        return True
    return False


def get_knowledge_base_name(kb_id: str) -> str:
    meta_path = os.path.join(KNOWLEDGE_DIR, kb_id, "meta.json")
    if os.path.exists(meta_path):
        try:
            with open(meta_path, "r", encoding="utf-8") as f:
                return json.load(f).get("name", kb_id)
        except Exception:
            pass
    return kb_id


# ── Document CRUD ──────────────────────────────────────────────────

def _list_files(kb_path: str) -> list[dict]:
    """列出知识库子目录中除 meta.json 外的文件"""
    files = []
    for fname in os.listdir(kb_path):
        if fname == "meta.json":
            continue
        fpath = os.path.join(kb_path, fname)
        if os.path.isfile(fpath):
            stat = os.stat(fpath)
            files.append({
                "name": fname,
                "size": stat.st_size,
                "modified": datetime.fromtimestamp(stat.st_mtime).isoformat(),
            })
    files.sort(key=lambda x: x["modified"], reverse=True)
    return files


def list_documents(kb_id: str) -> list[dict]:
    kb_path = os.path.join(KNOWLEDGE_DIR, kb_id)
    if not os.path.isdir(kb_path):
        return []
    return _list_files(kb_path)


def add_document(kb_id: str, src_path: str) -> bool:
    """将文件复制到知识库目录。返回是否成功。"""
    kb_path = os.path.join(KNOWLEDGE_DIR, kb_id)
    if not os.path.isdir(kb_path) or not os.path.isfile(src_path):
        return False
    dst = _unique_destination_path(kb_path, os.path.basename(src_path))
    shutil.copy2(src_path, dst)
    return True


def import_paths(kb_id: str, paths: list[str]) -> dict:
    kb_path = os.path.join(KNOWLEDGE_DIR, kb_id)
    result = {
        "imported": [],
        "skipped": [],
        "errors": [],
    }
    if not os.path.isdir(kb_path):
        result["errors"].append({"path": kb_id, "reason": "knowledge_base_not_found"})
        return result

    for raw_path in paths:
        if not raw_path:
            continue
        path = os.path.abspath(raw_path)
        if os.path.isfile(path):
            if not is_supported_knowledge_file(path):
                result["skipped"].append({"path": path, "reason": "unsupported"})
                continue
            try:
                if add_document(kb_id, path):
                    result["imported"].append(path)
                else:
                    result["errors"].append({"path": path, "reason": "copy_failed"})
            except Exception as e:
                result["errors"].append({"path": path, "reason": str(e)})
            continue

        if os.path.isdir(path):
            for root, _, filenames in os.walk(path):
                for filename in filenames:
                    candidate = os.path.join(root, filename)
                    if not is_supported_knowledge_file(candidate):
                        result["skipped"].append({"path": candidate, "reason": "unsupported"})
                        continue
                    try:
                        if add_document(kb_id, candidate):
                            result["imported"].append(candidate)
                        else:
                            result["errors"].append({"path": candidate, "reason": "copy_failed"})
                    except Exception as e:
                        result["errors"].append({"path": candidate, "reason": str(e)})
            continue

        result["skipped"].append({"path": path, "reason": "missing"})

    return result


def remove_document(kb_id: str, filename: str) -> bool:
    fpath = os.path.join(KNOWLEDGE_DIR, kb_id, filename)
    if os.path.isfile(fpath):
        os.remove(fpath)
        return True
    return False


# ── 文本提取（用于注入 LLM） ──────────────────────────────────────

_CSV_ENCODING_FALLBACKS = ("utf-8-sig", "utf-8", "gbk", "gb18030", "latin-1")


def _sniff_csv_encoding(file_path: str) -> str:
    """嗅探 CSV 文件编码：先用 chardet，再走 fallback 列表。"""
    try:
        with open(file_path, "rb") as f:
            head = f.read(65536)
    except OSError:
        return "utf-8"

    try:
        import chardet  # type: ignore
        guess = chardet.detect(head) or {}
        enc = (guess.get("encoding") or "").strip()
        conf = float(guess.get("confidence") or 0.0)
        if enc and conf >= 0.6:
            try:
                head.decode(enc)
                return enc
            except (LookupError, UnicodeDecodeError):
                pass
    except ImportError:
        pass

    for candidate in _CSV_ENCODING_FALLBACKS:
        try:
            head.decode(candidate)
            return candidate
        except (LookupError, UnicodeDecodeError):
            continue
    return "latin-1"


def _render_markdown_table(rows: list[list], max_rows: int = 200) -> str:
    if not rows:
        return ""
    header = [str(c) if c is not None else "" for c in rows[0]]
    out = ["| " + " | ".join(header) + " |",
           "| " + " | ".join("---" for _ in header) + " |"]
    body = rows[1:]
    shown = body[:max_rows]
    for row in shown:
        cells = []
        for cell in row:
            if cell is None:
                cells.append("")
            else:
                text = str(cell).replace("\n", " ").replace("|", "\\|")
                cells.append(text)
        # pad to header length
        while len(cells) < len(header):
            cells.append("")
        out.append("| " + " | ".join(cells) + " |")
    if len(body) > max_rows:
        out.append(f"... 省略 {len(body) - max_rows} 行 ...")
    return "\n".join(out)


def _extract_text(file_path: str) -> str:
    """从文件中提取纯文本内容"""
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == ".csv":
            enc = _sniff_csv_encoding(file_path)
            try:
                with open(file_path, "r", encoding=enc, errors="replace") as f:
                    body = f.read()
            except OSError as exc:
                return f"[CSV 读取失败: {os.path.basename(file_path)} — {exc}]"
            return f"# Encoding: {enc}\n{body}"

        if ext in (".txt", ".md", ".json", ".xml", ".html", ".log"):
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                return f.read()

        if ext == ".pdf":
            try:
                from PyPDF2 import PdfReader
                reader = PdfReader(file_path)
                return "\n".join(page.extract_text() or "" for page in reader.pages)
            except ImportError:
                return f"[PDF 文件无法解析，请安装 PyPDF2: {os.path.basename(file_path)}]"

        if ext == ".docx":
            try:
                import docx
                from docx.oxml.ns import qn  # type: ignore
                doc = docx.Document(file_path)
                parts: list[str] = []
                title = (doc.core_properties.title or "").strip() if doc.core_properties else ""
                if title:
                    parts.append(f"# {title}\n")

                # 用 element tree 顺序遍历 paragraphs 与 tables
                body = doc.element.body
                para_tag = qn("w:p")
                tbl_tag = qn("w:tbl")
                # Map para xml -> Paragraph object
                para_map = {p._element: p for p in doc.paragraphs}
                tbl_map = {t._element: t for t in doc.tables}

                for child in body.iterchildren():
                    if child.tag == para_tag:
                        para = para_map.get(child)
                        if para is None:
                            continue
                        text = (para.text or "").rstrip()
                        if not text:
                            continue
                        style = (para.style.name if para.style else "") or ""
                        if style.startswith("Heading"):
                            try:
                                level = int(style.split()[-1])
                            except (ValueError, IndexError):
                                level = 1
                            level = max(1, min(6, level))
                            parts.append("#" * level + " " + text)
                        else:
                            parts.append(text)
                    elif child.tag == tbl_tag:
                        tbl = tbl_map.get(child)
                        if tbl is None:
                            continue
                        rows = []
                        for row in tbl.rows:
                            rows.append([cell.text for cell in row.cells])
                        if rows:
                            parts.append(_render_markdown_table(rows))
                return "\n\n".join(parts)
            except ImportError:
                return f"[DOCX 文件无法解析，请安装 python-docx: {os.path.basename(file_path)}]"

        if ext == ".xlsx":
            try:
                import openpyxl
                wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
                try:
                    parts: list[str] = []
                    for ws in wb.worksheets:
                        rows = list(ws.iter_rows(values_only=True))
                        row_count = len(rows)
                        col_count = max((len(r) for r in rows), default=0)
                        parts.append(f"## Sheet: {ws.title} ({row_count} rows × {col_count} cols)")
                        if rows:
                            parts.append(_render_markdown_table([list(r) for r in rows]))
                    return "\n\n".join(parts)
                finally:
                    try:
                        wb.close()
                    except Exception:
                        pass
            except ImportError:
                return f"[XLSX 文件无法解析，请安装 openpyxl: {os.path.basename(file_path)}]"

        if ext == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                parts: list[str] = []
                for idx, slide in enumerate(prs.slides, start=1):
                    title_text = ""
                    body_lines: list[str] = []
                    tables_md: list[str] = []
                    notes_text = ""

                    title_shape = None
                    try:
                        title_shape = slide.shapes.title
                    except Exception:
                        title_shape = None

                    for shape in slide.shapes:
                        if shape is title_shape:
                            title_text = (shape.text or "").strip()
                            continue
                        if getattr(shape, "has_table", False):
                            tbl = shape.table
                            rows = []
                            for row in tbl.rows:
                                rows.append([cell.text for cell in row.cells])
                            if rows:
                                tables_md.append(_render_markdown_table(rows))
                            continue
                        if getattr(shape, "has_text_frame", False):
                            txt = (shape.text or "").strip()
                            if txt:
                                body_lines.append(txt)

                    try:
                        if slide.has_notes_slide:
                            notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
                    except Exception:
                        notes_text = ""

                    header = f"## Slide {idx}"
                    if title_text:
                        header = f"## Slide {idx}: {title_text}"
                    parts.append(header)
                    if body_lines:
                        parts.append("\n".join(body_lines))
                    for k, tmd in enumerate(tables_md, start=1):
                        parts.append(f"Table {k}:\n{tmd}")
                    if notes_text:
                        parts.append(f"Notes: {notes_text}")
                return "\n\n".join(parts)
            except ImportError:
                return f"[PPTX 文件无法解析，请安装 python-pptx: {os.path.basename(file_path)}]"

        # 尝试以纯文本打开未知类型
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()
    except Exception as e:
        return f"[提取失败: {os.path.basename(file_path)} — {e}]"


def load_knowledge_content(kb_id: str) -> str:
    """读取知识库全部文件的文本内容，拼接返回。"""
    kb_path = os.path.join(KNOWLEDGE_DIR, kb_id)
    if not os.path.isdir(kb_path):
        return ""
    parts = []
    for fname in sorted(os.listdir(kb_path)):
        if fname == "meta.json":
            continue
        fpath = os.path.join(kb_path, fname)
        if os.path.isfile(fpath):
            text = _extract_text(fpath)
            if text.strip():
                parts.append(f"--- {fname} ---\n{text}")
    return "\n\n".join(parts)


def _normalize_search_text(text: str) -> str:
    return str(text or "").casefold().strip()


def _tokenize_search_query(text: str) -> list[str]:
    normalized = _normalize_search_text(text)
    if not normalized:
        return []
    raw_tokens = re.findall(r"[a-z0-9_\-\.]{2,}|[\u4e00-\u9fff]{2,}", normalized)
    tokens = []
    for token in raw_tokens:
        tokens.append(token)
        if re.fullmatch(r"[\u4e00-\u9fff]{4,}", token):
            tokens.extend(token[index:index + 2] for index in range(len(token) - 1))
            tokens.extend(token[index:index + 3] for index in range(len(token) - 2))
    unique_tokens = []
    seen = set()
    for token in tokens:
        if token in seen:
            continue
        seen.add(token)
        unique_tokens.append(token)
    return unique_tokens


def _split_knowledge_chunks(text: str, max_chars: int = 1800) -> list[str]:
    cleaned = str(text or "").replace("\r\n", "\n").strip()
    if not cleaned:
        return []

    blocks = [block.strip() for block in re.split(r"\n(?=#{1,6}\s+)|\n{2,}", cleaned) if block.strip()]
    chunks = []
    for block in blocks:
        if len(block) <= max_chars:
            chunks.append(block)
            continue
        start = 0
        while start < len(block):
            chunk = block[start:start + max_chars].strip()
            if chunk:
                chunks.append(chunk)
            start += max_chars
    return chunks


def _score_knowledge_chunk(query_text: str, query_tokens: list[str], kb_name: str, filename: str, chunk: str) -> tuple[int, str]:
    name_text = _normalize_search_text(kb_name)
    file_text = _normalize_search_text(filename)
    chunk_text = _normalize_search_text(chunk)
    score = 0
    reasons = []

    if query_text:
        if query_text in name_text:
            score += 24
            reasons.append("完整请求命中知识库名")
        if query_text in file_text:
            score += 18
            reasons.append("完整请求命中文件名")
        if query_text in chunk_text:
            score += 12
            reasons.append("完整请求命中内容片段")

    matched = []
    for token in query_tokens:
        if token in name_text:
            score += 8
            matched.append(f"知识库名:{token}")
        elif token in file_text:
            score += 6
            matched.append(f"文件名:{token}")
        elif token in chunk_text:
            score += 2
            if len(matched) < 10:
                matched.append(f"内容:{token}")

    if matched:
        reasons.append("命中 " + ", ".join(matched[:10]))
    return score, "; ".join(reasons)


def search_knowledge_snippets(
    user_query: str,
    kb_ids: list[str] | None = None,
    max_kbs: int = 3,
    max_snippets: int = 6,
    snippet_chars: int = 1400,
) -> list[dict]:
    """按用户问题从知识库中检索最相关片段。

    当 kb_ids 为空时跨全部知识库自动匹配；当 kb_ids 非空时只在指定知识库内检索。
    返回 [{kb_id, kb_name, filename, score, reason, text}]。
    """
    query_text = _normalize_search_text(user_query)
    query_tokens = _tokenize_search_query(user_query)
    if not query_text or not query_tokens:
        return []

    allowed_ids = {str(kb_id) for kb_id in kb_ids or [] if str(kb_id).strip()}
    manual_scope = bool(allowed_ids)
    bases = list_knowledge_bases()
    if allowed_ids:
        bases = [kb for kb in bases if str(kb.get("id", "")) in allowed_ids]

    kb_matches = []
    for kb in bases:
        kb_id = str(kb.get("id", ""))
        kb_name = str(kb.get("name", kb_id))
        kb_path = os.path.join(KNOWLEDGE_DIR, kb_id)
        if not kb_id or not os.path.isdir(kb_path):
            continue

        scored_snippets = []
        fallback_snippets = []
        for fname in sorted(os.listdir(kb_path)):
            if fname == "meta.json":
                continue
            fpath = os.path.join(kb_path, fname)
            if not os.path.isfile(fpath):
                continue
            text = _extract_text(fpath)
            chunks = _split_knowledge_chunks(text)
            if chunks and len(fallback_snippets) < max_snippets:
                fallback_snippets.append({
                    "kb_id": kb_id,
                    "kb_name": kb_name,
                    "filename": fname,
                    "score": 1,
                    "reason": "手动选择知识库，未发现强关键词命中",
                    "text": chunks[0][:snippet_chars].strip(),
                })
            for chunk in chunks:
                score, reason = _score_knowledge_chunk(query_text, query_tokens, kb_name, fname, chunk)
                if score <= 0:
                    continue
                scored_snippets.append({
                    "kb_id": kb_id,
                    "kb_name": kb_name,
                    "filename": fname,
                    "score": score,
                    "reason": reason or "query terms matched this snippet",
                    "text": chunk[:snippet_chars].strip(),
                })

        scored_snippets.sort(key=lambda item: item["score"], reverse=True)
        if scored_snippets:
            selected = scored_snippets[:max_snippets]
            kb_score = sum(item["score"] for item in selected[:3])
            kb_matches.append((kb_score, selected))
        elif manual_scope and fallback_snippets:
            kb_matches.append((1, fallback_snippets[:max_snippets]))

    kb_matches.sort(key=lambda item: item[0], reverse=True)
    selected_snippets = []
    for _, snippets in kb_matches[:max_kbs]:
        selected_snippets.extend(snippets)
    selected_snippets.sort(key=lambda item: item["score"], reverse=True)
    return selected_snippets[:max_snippets]


def format_size(size_bytes: int) -> str:
    """人类可读的文件大小"""
    if size_bytes < 1024:
        return f"{size_bytes} B"
    if size_bytes < 1024 * 1024:
        return f"{size_bytes / 1024:.1f} KB"
    return f"{size_bytes / (1024 * 1024):.2f} MB"

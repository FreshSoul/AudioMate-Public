"""Unit tests for the structured document readers in AgentToolbox.

These tests construct ``AgentToolbox(None, None)`` directly — the new
``read_csv`` / ``read_xlsx`` / ``read_docx`` / ``read_pptx`` methods do
not require a Qt parent or a live WAAPI client. Each test writes a
small temporary fixture file and asserts on the returned structured
dict.
"""

from __future__ import annotations

import csv
import os
import sys
import tempfile

TEST_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_ROOT = os.path.abspath(os.path.join(TEST_DIR, "..", ".."))
if PROJECT_ROOT not in sys.path:
    sys.path.insert(0, PROJECT_ROOT)

from src.utils.agent_tools import AgentToolbox  # noqa: E402


def _toolbox() -> AgentToolbox:
    return AgentToolbox(parent_widget=None, waapi_client=None)


def _write_temp(content: bytes, suffix: str) -> str:
    fd, path = tempfile.mkstemp(suffix=suffix)
    os.close(fd)
    with open(path, "wb") as f:
        f.write(content)
    return path


# ---------------------------------------------------------------------------
# read_csv
# ---------------------------------------------------------------------------


def test_read_csv_utf8_with_header():
    body = "name,age\nAlice,30\nBob,25\nCarol,42\n".encode("utf-8")
    path = _write_temp(body, ".csv")
    try:
        result = _toolbox().read_csv(path)
        assert "error" not in result, result
        assert result["columns"] == ["name", "age"]
        assert len(result["rows"]) == 3
        assert result["truncated"] is False
        # dtype heuristic: age column entries are all digit strings → int
        assert result["dtypes"]["age"] == "int"
    finally:
        os.unlink(path)


def test_read_csv_gbk_chinese():
    body = "名称,数量\n苹果,3\n橙子,5\n".encode("gbk")
    path = _write_temp(body, ".csv")
    try:
        result = _toolbox().read_csv(path)
        assert "error" not in result, result
        # Encoding sniff should land on a GBK family
        assert result["encoding"].lower() in {"gbk", "gb2312", "gb18030"}
        assert result["columns"] == ["名称", "数量"]
        names = [row[0] for row in result["rows"]]
        assert "苹果" in names
        assert "橙子" in names
    finally:
        os.unlink(path)


def test_read_csv_semicolon_delimiter():
    body = "city;population\nShanghai;24000000\nBeijing;21000000\n".encode("utf-8")
    path = _write_temp(body, ".csv")
    try:
        result = _toolbox().read_csv(path)
        assert "error" not in result, result
        assert result["delimiter"] == ";"
        assert result["columns"] == ["city", "population"]
    finally:
        os.unlink(path)


# ---------------------------------------------------------------------------
# read_xlsx
# ---------------------------------------------------------------------------


def _make_xlsx_two_sheets() -> str:
    import openpyxl  # type: ignore
    wb = openpyxl.Workbook()
    ws1 = wb.active
    ws1.title = "Sheet1"
    ws1.append(["col_a", "col_b"])
    ws1.append([1, 2])
    ws1.append([3, 4])
    ws2 = wb.create_sheet("Data")
    ws2.append(["k", "v"])
    ws2.append(["x", "y"])
    fd, path = tempfile.mkstemp(suffix=".xlsx")
    os.close(fd)
    wb.save(path)
    return path


def test_read_xlsx_two_sheets():
    path = _make_xlsx_two_sheets()
    try:
        result = _toolbox().read_xlsx(path)
        assert "error" not in result, result
        names = [s["name"] for s in result["sheets"]]
        assert names == ["Sheet1", "Data"]
        # Sheet1 first row contains the column header strings.
        first = result["sheets"][0]
        assert first["columns"] == ["col_a", "col_b"]
        assert first["row_count"] == 3
    finally:
        os.unlink(path)


def test_read_xlsx_specific_sheet():
    path = _make_xlsx_two_sheets()
    try:
        result = _toolbox().read_xlsx(path, sheet="Data")
        assert "error" not in result, result
        assert len(result["sheets"]) == 1
        assert result["sheets"][0]["name"] == "Data"
        assert result["sheets"][0]["columns"] == ["k", "v"]
    finally:
        os.unlink(path)


def test_read_xlsx_preserves_numbers():
    import openpyxl  # type: ignore
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Numbers"
    ws.append(["a", "b"])
    ws.append([42, 3.14])
    fd, path = tempfile.mkstemp(suffix=".xlsx")
    os.close(fd)
    wb.save(path)
    try:
        result = _toolbox().read_xlsx(path)
        sheet = result["sheets"][0]
        # 2nd row is the data row: [42, 3.14] — values must keep numeric types
        data_row = sheet["rows"][1]
        assert isinstance(data_row[0], int)
        assert data_row[0] == 42
        assert isinstance(data_row[1], float)
        assert abs(data_row[1] - 3.14) < 1e-9
    finally:
        os.unlink(path)


# ---------------------------------------------------------------------------
# read_docx
# ---------------------------------------------------------------------------


def _make_docx_with_paragraph_and_table() -> str:
    import docx  # type: ignore
    doc = docx.Document()
    doc.add_paragraph("Hello, world.")
    table = doc.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "h1"
    table.rows[0].cells[1].text = "h2"
    table.rows[1].cells[0].text = "v1"
    table.rows[1].cells[1].text = "v2"
    fd, path = tempfile.mkstemp(suffix=".docx")
    os.close(fd)
    doc.save(path)
    return path


def test_read_docx_paragraphs_and_table():
    path = _make_docx_with_paragraph_and_table()
    try:
        result = _toolbox().read_docx(path)
        assert "error" not in result, result
        texts = [p["text"] for p in result["paragraphs"]]
        assert "Hello, world." in texts
        assert len(result["tables"]) == 1
        rows = result["tables"][0]["rows"]
        assert rows[0] == ["h1", "h2"]
        assert rows[1] == ["v1", "v2"]
    finally:
        os.unlink(path)


def test_read_docx_heading_levels():
    import docx  # type: ignore
    doc = docx.Document()
    doc.add_paragraph("Title One", style="Heading 1")
    doc.add_paragraph("Subtitle", style="Heading 2")
    doc.add_paragraph("Body text here.")
    fd, path = tempfile.mkstemp(suffix=".docx")
    os.close(fd)
    doc.save(path)
    try:
        result = _toolbox().read_docx(path)
        assert "error" not in result, result
        levels = [(h["level"], h["text"]) for h in result["headings"]]
        assert (1, "Title One") in levels
        assert (2, "Subtitle") in levels
    finally:
        os.unlink(path)


# ---------------------------------------------------------------------------
# read_pptx
# ---------------------------------------------------------------------------


def _make_pptx_two_slides_with_notes() -> str:
    from pptx import Presentation  # type: ignore
    prs = Presentation()
    layout = prs.slide_layouts[1]  # title + content
    s1 = prs.slides.add_slide(layout)
    s1.shapes.title.text = "First slide"
    # body placeholder is usually index 1
    for ph in s1.placeholders:
        if ph.placeholder_format.idx == 1:
            ph.text = "first body"
            break
    s1.notes_slide.notes_text_frame.text = "speaker note one"

    s2 = prs.slides.add_slide(layout)
    s2.shapes.title.text = "Second slide"
    for ph in s2.placeholders:
        if ph.placeholder_format.idx == 1:
            ph.text = "second body"
            break
    s2.notes_slide.notes_text_frame.text = "speaker note two"

    fd, path = tempfile.mkstemp(suffix=".pptx")
    os.close(fd)
    prs.save(path)
    return path


def test_read_pptx_title_body_notes():
    path = _make_pptx_two_slides_with_notes()
    try:
        result = _toolbox().read_pptx(path)
        assert "error" not in result, result
        assert result["slide_count"] == 2
        assert result["slides"][0]["title"] == "First slide"
        assert "first body" in result["slides"][0]["body_text"]
        assert result["slides"][0]["notes"] == "speaker note one"
        assert result["slides"][1]["title"] == "Second slide"
        assert result["slides"][1]["notes"] == "speaker note two"
    finally:
        os.unlink(path)


def test_read_pptx_slide_table():
    from pptx import Presentation  # type: ignore
    from pptx.util import Inches  # type: ignore
    prs = Presentation()
    layout = prs.slide_layouts[5]  # blank-ish with title
    slide = prs.slides.add_slide(layout)
    try:
        slide.shapes.title.text = "Table slide"
    except AttributeError:
        pass
    rows, cols = 2, 2
    left = top = Inches(1)
    width = Inches(4)
    height = Inches(1.5)
    shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = shape.table
    table.cell(0, 0).text = "A"
    table.cell(0, 1).text = "B"
    table.cell(1, 0).text = "1"
    table.cell(1, 1).text = "2"

    fd, path = tempfile.mkstemp(suffix=".pptx")
    os.close(fd)
    prs.save(path)
    try:
        result = _toolbox().read_pptx(path)
        assert "error" not in result, result
        tables = result["slides"][0]["tables"]
        assert len(tables) == 1
        assert tables[0]["rows"][0] == ["A", "B"]
        assert tables[0]["rows"][1] == ["1", "2"]
    finally:
        os.unlink(path)


if __name__ == "__main__":  # pragma: no cover — manual run helper
    fns = [v for k, v in dict(globals()).items() if k.startswith("test_")]
    failures = 0
    for fn in fns:
        try:
            fn()
            print(f"PASS {fn.__name__}")
        except Exception as exc:
            failures += 1
            print(f"FAIL {fn.__name__}: {type(exc).__name__}: {exc}")
    if failures:
        print(f"\n{failures}/{len(fns)} failed")
        sys.exit(1)
    print(f"\nOK ({len(fns)} tests passed)")
